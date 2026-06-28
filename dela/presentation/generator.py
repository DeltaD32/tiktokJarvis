"""Slide generator — build a .pptx deck from a storyline using a stored style.

The generator takes a storyline (list of slide specs) and a style slug, loads
the style profile (colors, fonts, geometry), opens the style's source.pptx as
the template base, and builds each slide using pptx_lib helpers.

Storyline format (list of dicts):
  [
    {
      "layout": 7,           # layout index in the template
      "title": "Slide Title",
      "content": ["bullet 1", "bullet 2"],   # for content layouts
      "layout_type": "bullets",  # "bullets", "title_only", "hero_number", etc.
      "hero_number": "2.3B",     # for hero_number layout
      "hero_subtitle": "Revenue",
      "notes": "Speaker notes text",
    },
    ...
  ]

The generator supports simple layout types directly. For complex compositions
(free-area layouts with cards, pillars, timelines, etc.), the caller provides
a "build_code" string — Python code that runs with pptx_lib in scope, similar
to the opencode-galaxy sub-agent pattern. This is executed via the sandbox.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from dela.presentation.pptx_lib.constants import load_style_constants
from dela.presentation.pptx_lib import (
    COLOR_PRIMARY,
    COLOR_ACCENT,
    DARK,
    WHITE,
    FREE_LEFT,
    FREE_TOP,
    FREE_WIDTH,
    FREE_HEIGHT,
    add_textbox,
    add_bullet_paragraph,
    add_hero_number,
    add_pillars,
    add_mece_tiles,
    add_styled_table,
    add_chevron_chain,
    add_card_with_accent,
    send_to_back,
    remove_empty_placeholders,
    set_notes,
    suppress_bullet,
)
from dela.presentation.style_registry import style_template_path, load_style_profile


def generate(
    style_slug: str,
    storyline: list[dict[str, Any]],
    output_path: str,
    footer: str = "",
    confidentiality: str = "",
) -> str:
    """Generate a .pptx file from a storyline using a stored style.

    Returns the output path on success, or an error message.
    """
    # Load the style template
    template_path = style_template_path(style_slug)
    if not template_path.exists():
        return f"Style '{style_slug}' has no source.pptx template. Clone it first."

    # Load style constants (colors, fonts, geometry)
    load_style_constants(style_slug)

    # Open the template
    prs = Presentation(str(template_path))

    # Delete existing slides from the template
    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Update footer on slide master if shapes exist
    if footer:
        _update_master_footer(prs, footer, confidentiality)

    # Build each slide
    for i, spec in enumerate(storyline):
        _build_slide(prs, spec, i + 1)

    # Save
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return str(output)


def _update_master_footer(prs: Presentation, footer: str, confidentiality: str) -> None:
    """Update footer text on the slide master, if footer shapes exist."""
    master = prs.slide_masters[0]
    for shape in master.shapes:
        if not shape.has_text_frame:
            continue
        # Generic footer shape detection by name (works with any template)
        name = shape.name.lower()
        if "fußzeile" in name or "footer" in name or "fuzeile" in name:
            _set_para_text(shape.text_frame.paragraphs[0], footer)
        elif "confidential" in name or "textfeld 3" in name.lower():
            if confidentiality:
                _set_para_text(shape.text_frame.paragraphs[0], confidentiality)


def _set_para_text(para, text: str) -> None:
    """Set paragraph text without destroying run-level formatting."""
    runs = para.runs
    if runs:
        runs[0].text = text
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        para.text = text


def _build_slide(prs: Presentation, spec: dict[str, Any], slide_num: int) -> None:
    """Build a single slide from a storyline spec."""
    layout_idx = spec.get("layout", 12)  # default to free-area layout
    layout_type = spec.get("layout_type", "bullets")
    title = spec.get("title", "")
    notes = spec.get("notes", "")

    # Clamp layout index to available layouts
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = len(prs.slide_layouts) - 1

    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Set title (placeholder 0 on most templates)
    if title:
        try:
            title_ph = slide.placeholders[0]
            title_ph.text = title
            # Suppress inherited bullets on title
            for p in title_ph.text_frame.paragraphs:
                suppress_bullet(p)
        except (KeyError, IndexError):
            # No title placeholder — add a textbox
            add_textbox(slide, FREE_LEFT, Cm(0.5), FREE_WIDTH, Cm(1.5),
                        text=title, font_size=Pt(28), bold=True, color=COLOR_PRIMARY)

    # Build content based on layout type
    builder = _LAYOUT_BUILDERS.get(layout_type, _build_bullets)
    builder(slide, spec)

    # Clean up empty placeholders
    remove_empty_placeholders(slide)

    # Set speaker notes
    if notes:
        set_notes(slide, notes)


def _build_bullets(slide, spec: dict[str, Any]) -> None:
    """Standard bullet content in the body placeholder."""
    content = spec.get("content", [])
    if not content:
        return
    try:
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        tf.clear()
        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK
            p.space_before = Pt(6)
    except (KeyError, IndexError):
        # No body placeholder — add a textbox
        add_textbox(slide, FREE_LEFT, FREE_TOP, FREE_WIDTH, FREE_HEIGHT,
                    text="\n".join(f"• {item}" for item in content),
                    font_size=Pt(18), color=DARK)


def _build_title_only(slide, spec: dict[str, Any]) -> None:
    """Title only — no body content."""
    pass  # Title already set in _build_slide


def _build_hero_number(slide, spec: dict[str, Any]) -> None:
    """Large centered number with subtitle."""
    number = spec.get("hero_number", "")
    subtitle = spec.get("hero_subtitle", "")
    add_hero_number(slide, number, subtitle)


def _build_pillars(slide, spec: dict[str, Any]) -> None:
    """Pillars layout — columns with headers and bullets."""
    pillars = spec.get("pillars", [])
    if not pillars:
        return
    add_pillars(slide, pillars, FREE_LEFT, FREE_TOP, FREE_WIDTH, FREE_HEIGHT)


def _build_mece_tiles(slide, spec: dict[str, Any]) -> None:
    """MECE tiles — equal-width vertical tiles."""
    tiles = spec.get("tiles", [])
    if not tiles:
        return
    add_mece_tiles(slide, tiles, FREE_LEFT, FREE_TOP, FREE_WIDTH, FREE_HEIGHT)


def _build_table(slide, spec: dict[str, Any]) -> None:
    """Styled table."""
    headers = spec.get("table_headers", [])
    rows = spec.get("table_rows", [])
    if not headers:
        return
    add_styled_table(slide, headers, rows, FREE_LEFT, FREE_TOP, FREE_WIDTH, Cm(10))


def _build_chevron(slide, spec: dict[str, Any]) -> None:
    """Chevron chain — horizontal arrows with numbered badges."""
    items = spec.get("chevron_items", [])
    if not items:
        return
    add_chevron_chain(slide, items)


def _build_cards(slide, spec: dict[str, Any]) -> None:
    """Cards with accent strips."""
    cards = spec.get("cards", [])
    if not cards:
        return
    n = len(cards)
    from pptx.util import Emu
    col_w = int(FREE_WIDTH) // n
    for i, card_spec in enumerate(cards):
        left = int(FREE_LEFT) + i * col_w
        title = card_spec.get("title", "")
        body = card_spec.get("body", "")
        accent = card_spec.get("accent_color", COLOR_ACCENT)
        add_card_with_accent(slide, Emu(left), FREE_TOP, Emu(col_w - int(Cm(0.4))), FREE_HEIGHT, accent)
        add_textbox(slide, Emu(left + int(Cm(0.4))), FREE_TOP + Cm(0.4),
                    Emu(col_w - int(Cm(0.8))), Cm(1.5),
                    text=title, font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
        add_textbox(slide, Emu(left + int(Cm(0.4))), FREE_TOP + Cm(2),
                    Emu(col_w - int(Cm(0.8))), FREE_HEIGHT - Cm(2.5),
                    text=body, font_size=Pt(14), color=DARK)


def _build_key_message(slide, spec: dict[str, Any]) -> None:
    """Single large key message centered on the slide."""
    message = spec.get("key_message", "")
    if not message:
        return
    add_textbox(slide, FREE_LEFT, FREE_TOP, FREE_WIDTH, FREE_HEIGHT,
                text=message, font_size=Pt(36), bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)


# Registry of layout type → builder function
_LAYOUT_BUILDERS: dict[str, Any] = {
    "bullets": _build_bullets,
    "title_only": _build_title_only,
    "hero_number": _build_hero_number,
    "pillars": _build_pillars,
    "mece_tiles": _build_mece_tiles,
    "table": _build_table,
    "chevron": _build_chevron,
    "cards": _build_cards,
    "key_message": _build_key_message,
}