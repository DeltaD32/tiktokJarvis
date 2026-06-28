#!/usr/bin/env python3
"""
ppt-style-cloner — extract and persist the full visual style of any .pptx file
into the local style registry at ~/.config/opencode/ppt-styles/.

Usage:
  python clone_style.py <path/to/source.pptx> [--name "My Style"] [--description "..."]
  python clone_style.py <path/to/source.pptx> --overwrite   # replace existing slug

Outputs:
  ~/.config/opencode/ppt-styles/<slug>/
    ├── style.json              Full extracted style profile
    ├── brand-guide.md          Human-readable brand guide (usable by agents)
    ├── title-bg.jpeg           Extracted title slide background (if present)
    └── source.pptx             Copy of the original file (for re-extraction)

  ~/.config/opencode/ppt-styles/registry.json   Updated index

Extraction coverage
  ✓ Theme color scheme  — all 12 named slots (dk1/lt1/dk2/lt2/accent1-6/hlink/folHlink)
                          with srgbClr, sysClr, and scRgbClr variants
  ✓ Theme font scheme   — majorFont (+mj-lt) and minorFont (+mn-lt) resolved to real names
  ✓ Master txStyles     — titleStyle + bodyStyle: font size, color, cap, bullet char/font/color,
                          line-spacing, space-before/after, indent, alignment per level (1-9)
  ✓ Master background   — solid, gradient, bgRef (with scheme resolution), or picture
  ✓ Layout backgrounds  — per-layout bg fills; solid/gradient/picture/inherited; dark/light flag
  ✓ Layout title colors — per-layout title placeholder font color (scheme or hex)
  ✓ Layout placeholders — idx, type, name, position/size in inches
  ✓ Master named shapes — FußzeileAU1, SeitenzahlAU1, BMW Group Trennlinie etc.
                          with position/size, fill color, line color, text content
  ✓ Slide shape fills   — every shape across every slide: solid fill (hex or scheme-resolved),
                          gradient stops, picture fill flag
  ✓ Typography heuristic — dominant font, title/body/footer sizes from combined corpus
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import sys
from collections import Counter, OrderedDict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

# ── required dep ─────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERROR: python-pptx is not installed.  Run:  pip install python-pptx")
    sys.exit(1)

# lxml is bundled with python-pptx, so it's always available
from lxml import etree

# ── namespace map ─────────────────────────────────────────────────────────────
A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
P   = "http://schemas.openxmlformats.org/presentationml/2006/main"
R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

def _a(tag: str) -> str:  return f"{{{A}}}{tag}"
def _p(tag: str) -> str:  return f"{{{P}}}{tag}"

# ── constants ─────────────────────────────────────────────────────────────────
REGISTRY_DIR  = Path(__file__).resolve().parent.parent.parent / "dela_state" / "styles"
REGISTRY_FILE = REGISTRY_DIR / "registry.json"

# Ordered list of the 12 clrScheme slot names
CLR_SCHEME_SLOTS = ["dk1","lt1","dk2","lt2",
                    "accent1","accent2","accent3","accent4","accent5","accent6",
                    "hlink","folHlink"]

# Scheme color → semantic label
CLR_LABELS = {
    "dk1": "Text/Dark 1 (primary text)",
    "lt1": "Background/Light 1 (slide bg)",
    "dk2": "Dark 2 (primary brand)",
    "lt2": "Light 2",
    "accent1": "Accent 1",
    "accent2": "Accent 2",
    "accent3": "Accent 3",
    "accent4": "Accent 4",
    "accent5": "Accent 5",
    "accent6": "Accent 6",
    "hlink":    "Hyperlink",
    "folHlink": "Followed Hyperlink",
    # aliases used inside slides
    "tx1": "Text 1 → alias for dk1",
    "tx2": "Text 2 → alias for dk2",
    "bg1": "Background 1 → alias for lt1",
    "bg2": "Background 2 → alias for lt2",
    "phClr": "Placeholder Color (context-dependent)",
}


# ═════════════════════════════════════════════════════════════════════════════
# LOW-LEVEL XML HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def _srgb(el: etree._Element | None) -> str | None:
    """Extract hex from <a:srgbClr> child, or None."""
    if el is None:
        return None
    child = el.find(_a("srgbClr"))
    if child is not None:
        val = child.get("val")
        return f"#{val.upper()}" if val else None
    # sysClr has lastClr attribute
    child = el.find(_a("sysClr"))
    if child is not None:
        val = child.get("lastClr")
        return f"#{val.upper()}" if val else None
    # scRgbClr (percentage, 0–100000)
    child = el.find(_a("scRgbClr"))
    if child is not None:
        try:
            r = round(int(child.get("r", "0")) * 255 / 100000)
            g = round(int(child.get("g", "0")) * 255 / 100000)
            b = round(int(child.get("b", "0")) * 255 / 100000)
            return f"#{r:02X}{g:02X}{b:02X}"
        except Exception:
            pass
    return None


def _scheme_val(el: etree._Element | None) -> str | None:
    """Extract the val attr of <a:schemeClr>, or None."""
    if el is None:
        return None
    child = el.find(_a("schemeClr"))
    return child.get("val") if child is not None else None


def _safe_pt(sz_str: str | None) -> float | None:
    """Convert hundredths-of-point string (e.g. '2600') → pt float."""
    try:
        return round(int(sz_str) / 100, 1)
    except (TypeError, ValueError):
        return None


def _safe_emu_in(emu: Any) -> float | None:
    try:
        return round(float(emu) / 914400, 3)
    except Exception:
        return None


def _spc_pts(el: etree._Element | None) -> float | None:
    """Extract point value from <a:spcPts> inside <a:spcBef>/<a:spcAft>."""
    if el is None:
        return None
    pts = el.find(_a("spcPts"))
    if pts is not None:
        try:
            return round(int(pts.get("val", "0")) / 100, 1)
        except Exception:
            pass
    pct = el.find(_a("spcPct"))
    if pct is not None:
        try:
            return f"{int(pct.get('val', '0')) / 1000:.0f}%"
        except Exception:
            pass
    return None


# ═════════════════════════════════════════════════════════════════════════════
# THEME PART HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def _get_theme_part(master) -> etree._Element | None:
    """Return the parsed lxml element for the theme part of a slide master."""
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            try:
                return etree.fromstring(rel._target._blob)
            except Exception:
                pass
    return None


# ═════════════════════════════════════════════════════════════════════════════
# 1. THEME COLOR SCHEME
# ═════════════════════════════════════════════════════════════════════════════

def extract_theme_colors(master) -> dict:
    """
    Parse <a:clrScheme> from the theme part.
    Returns an OrderedDict of slot_name → {"hex": "#RRGGBB", "label": "..."}
    Also builds a resolver dict for scheme alias → hex (used elsewhere).
    """
    result: dict[str, dict] = {}
    tree = _get_theme_part(master)
    if tree is None:
        return result

    clr_scheme = tree.find(f".//{_a('clrScheme')}")
    if clr_scheme is None:
        return result

    for slot in CLR_SCHEME_SLOTS:
        slot_el = clr_scheme.find(_a(slot))
        if slot_el is None:
            continue
        hex_val = _srgb(slot_el)
        result[slot] = {
            "hex": hex_val,
            "label": CLR_LABELS.get(slot, slot),
            "scheme_name": clr_scheme.get("name", ""),
        }

    return result


def build_scheme_resolver(theme_colors: dict) -> dict[str, str]:
    """
    Returns a dict mapping every scheme alias → resolved hex string.
    e.g. {"dk1": "#000000", "tx1": "#000000", "bg1": "#FFFFFF", ...}
    """
    r: dict[str, str] = {}
    for slot, info in theme_colors.items():
        if info.get("hex"):
            r[slot] = info["hex"]
    # aliases
    alias_map = {"tx1": "dk1", "tx2": "dk2", "bg1": "lt1", "bg2": "lt2"}
    for alias, canonical in alias_map.items():
        if canonical in r:
            r[alias] = r[canonical]
    return r


def resolve_color(el: etree._Element | None, resolver: dict[str, str]) -> dict | None:
    """
    Given a parent element that may contain solidFill/gradFill/blipFill/noFill,
    return a compact fill descriptor.
    """
    if el is None:
        return None
    solid = el.find(_a("solidFill"))
    if solid is not None:
        hex_c = _srgb(solid)
        scheme = _scheme_val(solid)
        if hex_c is None and scheme and scheme in resolver:
            hex_c = resolver[scheme]
        return {"type": "solid", "hex": hex_c, "scheme": scheme}

    grad = el.find(_a("gradFill"))
    if grad is not None:
        stops = []
        for gs in grad.findall(_a("gs")):
            pos = gs.get("pos")
            hex_c = _srgb(gs)
            scheme = _scheme_val(gs)
            if hex_c is None and scheme and scheme in resolver:
                hex_c = resolver[scheme]
            stops.append({"pos": int(pos)/1000 if pos else None,
                           "hex": hex_c, "scheme": scheme})
        return {"type": "gradient", "stops": stops}

    if el.find(_a("blipFill")) is not None:
        return {"type": "picture"}

    if el.find(_a("noFill")) is not None:
        return {"type": "none"}

    return None


# ═════════════════════════════════════════════════════════════════════════════
# 2. THEME FONT SCHEME
# ═════════════════════════════════════════════════════════════════════════════

def extract_theme_fonts(master) -> dict:
    """
    Parse <a:fontScheme> from the theme part.
    Returns {"scheme_name": ..., "major_latin": ..., "minor_latin": ...,
             "major_ea": ..., "minor_ea": ...,
             "resolved": {"+mj-lt": "<real name>", "+mn-lt": "<real name>"}}
    """
    tree = _get_theme_part(master)
    if tree is None:
        return {}

    font_scheme = tree.find(f".//{_a('fontScheme')}")
    if font_scheme is None:
        return {}

    def _typeface(parent_tag: str, child_tag: str) -> str | None:
        parent = font_scheme.find(_a(parent_tag))
        if parent is None:
            return None
        child = parent.find(_a(child_tag))
        return child.get("typeface") if child is not None else None

    major_latin = _typeface("majorFont", "latin")
    minor_latin = _typeface("minorFont", "latin")
    major_ea    = _typeface("majorFont", "ea")
    minor_ea    = _typeface("minorFont", "ea")

    return {
        "scheme_name": font_scheme.get("name", ""),
        "major_latin": major_latin,   # resolves +mj-lt
        "minor_latin": minor_latin,   # resolves +mn-lt
        "major_ea": major_ea,
        "minor_ea": minor_ea,
        "resolved": {
            "+mj-lt": major_latin or "Arial",
            "+mn-lt": minor_latin or "Arial",
            "+mj-ea": major_ea or "",
            "+mn-ea": minor_ea or "",
        },
    }


def resolve_typeface(typeface: str | None, font_resolver: dict[str, str]) -> str | None:
    """Replace +mj-lt / +mn-lt tokens with their real font name."""
    if typeface is None:
        return None
    return font_resolver.get(typeface, typeface)


# ═════════════════════════════════════════════════════════════════════════════
# 3. MASTER txStyles  (titleStyle + bodyStyle)
# ═════════════════════════════════════════════════════════════════════════════

def _parse_pPr(pPr: etree._Element,
               scheme_resolver: dict[str, str],
               font_resolver: dict[str, str]) -> dict:
    """Extract everything useful from a <a:lvlNpPr> element."""
    out: dict = {}

    # Alignment
    algn = pPr.get("algn")
    if algn:
        out["align"] = algn

    # Indent / margin
    mar_l  = pPr.get("marL")
    indent = pPr.get("indent")
    if mar_l:  out["margin_left_emu"] = int(mar_l)
    if indent: out["indent_emu"]      = int(indent)

    # Spacing
    spc_bef = pPr.find(_a("spcBef"))
    spc_aft = pPr.find(_a("spcAft"))
    lnSpc   = pPr.find(_a("lnSpc"))
    if spc_bef is not None: out["space_before_pt"]  = _spc_pts(spc_bef)
    if spc_aft is not None: out["space_after_pt"]   = _spc_pts(spc_aft)
    if lnSpc   is not None: out["line_spacing"]     = _spc_pts(lnSpc)

    # Bullet
    bu_none   = pPr.find(_a("buNone"))
    bu_char   = pPr.find(_a("buChar"))
    bu_auto   = pPr.find(_a("buAutoNum"))
    bu_font   = pPr.find(_a("buFont"))
    bu_clr    = pPr.find(_a("buClr"))
    bu_size   = pPr.find(_a("buSzPct"))

    if bu_none is not None:
        out["bullet"] = {"type": "none"}
    elif bu_char is not None:
        char = bu_char.get("char", "")
        font_name = bu_font.get("typeface") if bu_font is not None else None
        font_name = resolve_typeface(font_name, font_resolver)
        bullet_color_hex = None
        if bu_clr is not None:
            hex_c = _srgb(bu_clr)
            scheme = _scheme_val(bu_clr)
            if hex_c is None and scheme and scheme in scheme_resolver:
                hex_c = scheme_resolver[scheme]
            bullet_color_hex = hex_c
            bullet_color_scheme = scheme
        out["bullet"] = {
            "type": "char",
            "char": char,
            "font": font_name,
            "color_hex": bullet_color_hex,
            "color_scheme": _scheme_val(bu_clr) if bu_clr is not None else None,
        }
    elif bu_auto is not None:
        out["bullet"] = {"type": "auto", "scheme": bu_auto.get("type")}

    # Default run properties
    def_rpr = pPr.find(_a("defRPr"))
    if def_rpr is not None:
        rpr_out: dict = {}
        sz = def_rpr.get("sz")
        if sz: rpr_out["size_pt"] = _safe_pt(sz)
        bold   = def_rpr.get("b")
        italic = def_rpr.get("i")
        cap    = def_rpr.get("cap")
        kern   = def_rpr.get("kern")
        if bold   is not None: rpr_out["bold"]   = bold   == "1"
        if italic is not None: rpr_out["italic"] = italic == "1"
        if cap    is not None: rpr_out["cap"]    = cap     # "all", "small", "none"

        # Font
        latin = def_rpr.find(_a("latin"))
        if latin is not None:
            tf = latin.get("typeface")
            rpr_out["font"] = resolve_typeface(tf, font_resolver)

        # Text color
        solid_fill = def_rpr.find(_a("solidFill"))
        if solid_fill is not None:
            hex_c  = _srgb(solid_fill)
            scheme = _scheme_val(solid_fill)
            if hex_c is None and scheme and scheme in scheme_resolver:
                hex_c = scheme_resolver[scheme]
            rpr_out["color_hex"]    = hex_c
            rpr_out["color_scheme"] = scheme

        if rpr_out:
            out["default_run_props"] = rpr_out

    return out


def extract_tx_styles(master,
                      scheme_resolver: dict[str, str],
                      font_resolver: dict[str, str]) -> dict:
    """
    Parse <p:txStyles> from the slide master element.
    Returns:
      {
        "title": {1: {...pPr fields...}},
        "body":  {1: {...}, 2: {...}, ..., 9: {...}},
      }
    """
    result: dict[str, dict] = {"title": {}, "body": {}}

    tx_styles = master.element.find(_p("txStyles"))
    if tx_styles is None:
        return result

    for style_tag, key in [("titleStyle", "title"), ("bodyStyle", "body")]:
        style_el = tx_styles.find(_p(style_tag))
        if style_el is None:
            continue
        for lvl in range(1, 10):
            lvl_el = style_el.find(_a(f"lvl{lvl}pPr"))
            if lvl_el is None:
                continue
            result[key][lvl] = _parse_pPr(lvl_el, scheme_resolver, font_resolver)

    return result


# ═════════════════════════════════════════════════════════════════════════════
# 4. MASTER BACKGROUND
# ═════════════════════════════════════════════════════════════════════════════

def _parse_bg_element(bg_el: etree._Element | None,
                      scheme_resolver: dict[str, str]) -> dict:
    """Parse a <p:bg> element into a fill descriptor."""
    if bg_el is None:
        return {"type": "inherited"}

    # p:bgRef (theme background reference)
    bg_ref = bg_el.find(_p("bgRef"))
    if bg_ref is not None:
        idx = bg_ref.get("idx")
        scheme = _scheme_val(bg_ref)
        hex_c  = _srgb(bg_ref)
        if hex_c is None and scheme and scheme in scheme_resolver:
            hex_c = scheme_resolver[scheme]
        return {"type": "bg_ref", "idx": idx, "scheme": scheme, "hex": hex_c}

    # p:bgPr
    bg_pr = bg_el.find(_p("bgPr"))
    if bg_pr is None:
        return {"type": "inherited"}

    solid = bg_pr.find(_a("solidFill"))
    if solid is not None:
        hex_c  = _srgb(solid)
        scheme = _scheme_val(solid)
        if hex_c is None and scheme and scheme in scheme_resolver:
            hex_c = scheme_resolver[scheme]
        return {"type": "solid", "hex": hex_c, "scheme": scheme}

    grad = bg_pr.find(_a("gradFill"))
    if grad is not None:
        stops = []
        for gs in grad.findall(f".//{_a('gs')}"):
            pos   = gs.get("pos")
            hex_c = _srgb(gs)
            sc    = _scheme_val(gs)
            if hex_c is None and sc and sc in scheme_resolver:
                hex_c = scheme_resolver[sc]
            stops.append({"pos": int(pos)/1000 if pos else None,
                           "hex": hex_c, "scheme": sc})
        return {"type": "gradient", "stops": stops}

    blip = bg_pr.find(_a("blipFill"))
    if blip is not None:
        return {"type": "picture"}

    return {"type": "unknown"}


def extract_master_background(master, scheme_resolver: dict[str, str]) -> dict:
    """Extract the background fill of the slide master."""
    bg_el = master.element.find(_p("bg"))
    return _parse_bg_element(bg_el, scheme_resolver)


# ═════════════════════════════════════════════════════════════════════════════
# 5. LAYOUT BACKGROUNDS + TITLE COLORS
# ═════════════════════════════════════════════════════════════════════════════

def _is_dark_fill(fill: dict, scheme_resolver: dict[str, str]) -> bool | None:
    """
    Heuristic: is this fill dark?
    Returns True, False, or None (unknown).
    """
    if fill.get("type") == "solid":
        hex_c = fill.get("hex")
        if hex_c and hex_c.startswith("#"):
            try:
                r = int(hex_c[1:3], 16)
                g = int(hex_c[3:5], 16)
                b = int(hex_c[5:7], 16)
                luminance = 0.299*r + 0.587*g + 0.114*b
                return luminance < 128
            except Exception:
                pass
        scheme = fill.get("scheme")
        if scheme in ("dk1", "dk2", "tx1", "tx2"):
            return True
        if scheme in ("lt1", "lt2", "bg1", "bg2"):
            return False
    if fill.get("type") == "gradient":
        for stop in fill.get("stops", []):
            dark = _is_dark_fill({"type": "solid",
                                   "hex": stop.get("hex"),
                                   "scheme": stop.get("scheme")},
                                  scheme_resolver)
            if dark is not None:
                return dark
    return None


# Keywords in layout names that signal a dark or gradient background
_DARK_NAME_TOKENS     = {"dark", "dunkel", "night", "nacht", "black", "schwarz"}
_GRADIENT_NAME_TOKENS = {"gradient", "verlauf", "genai", "aiconic"}
_LIGHT_NAME_TOKENS    = {"white", "weiss", "weiß", "hell", "light"}

def _is_dark_by_name(layout_name: str) -> bool | None:
    """
    Fallback: infer dark/light from the layout name when the background fill
    is marked 'inherited' (i.e. the deck uses picture-placeholder backgrounds
    or shape-layer backgrounds rather than <p:bgPr>).
    Returns True (dark), False (light), or None (unknown).
    """
    lower = layout_name.lower()
    tokens = set(lower.replace("|", " ").replace("_", " ").split())
    if tokens & _DARK_NAME_TOKENS:
        return True
    if tokens & _GRADIENT_NAME_TOKENS:
        return True   # gradient layouts typically have non-white backgrounds
    if tokens & _LIGHT_NAME_TOKENS:
        return False
    return None


def extract_layout_info(master,
                        scheme_resolver: dict[str, str],
                        font_resolver: dict[str, str],
                        master_is_dark: bool | None = None) -> list[dict]:
    """
    For every layout: name, background fill, is_dark flag, title color,
    and placeholder details (idx, type, name, position, size).
    master_is_dark is used as final fallback when both fill and name detection
    return None (i.e. layout inherits background from master).
    """
    layouts = []
    for layout in master.slide_layouts:
        # Background
        bg_el  = layout.element.find(_p("bg"))
        bg     = _parse_bg_element(bg_el, scheme_resolver)
        is_dark = _is_dark_fill(bg, scheme_resolver)
        if bg["type"] == "inherited":
            is_dark = None  # resolved from master at render time

        # Title placeholder color
        title_color: dict | None = None
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0:
                # Look for color in defRPr or rPr inside txBody
                for rpr_tag in [_a("defRPr"), _a("rPr")]:
                    for rpr in ph.element.iter(rpr_tag):
                        sf = rpr.find(_a("solidFill"))
                        if sf is not None:
                            hex_c  = _srgb(sf)
                            scheme = _scheme_val(sf)
                            if hex_c is None and scheme and scheme in scheme_resolver:
                                hex_c = scheme_resolver[scheme]
                            title_color = {"hex": hex_c, "scheme": scheme}
                            break
                    if title_color:
                        break
                # Also check lvl1pPr inside txBody lstStyle
                if title_color is None:
                    for lvl1 in ph.element.iter(_a("lvl1pPr")):
                        for rpr in lvl1.iter(_a("defRPr")):
                            sf = rpr.find(_a("solidFill"))
                            if sf is not None:
                                hex_c  = _srgb(sf)
                                scheme = _scheme_val(sf)
                                if hex_c is None and scheme and scheme in scheme_resolver:
                                    hex_c = scheme_resolver[scheme]
                                title_color = {"hex": hex_c, "scheme": scheme}
                                break
                break

        # Placeholders
        ph_list = []
        for ph in layout.placeholders:
            ph_list.append({
                "idx":      ph.placeholder_format.idx,
                "type":     str(ph.placeholder_format.type).split(".")[-1],
                "name":     ph.name,
                "left_in":  _safe_emu_in(ph.left),
                "top_in":   _safe_emu_in(ph.top),
                "width_in": _safe_emu_in(ph.width),
                "height_in":_safe_emu_in(ph.height),
            })

        # If fill-based detection gave None, fall back to layout name heuristic
        if is_dark is None:
            is_dark = _is_dark_by_name(layout.name)
        # Final fallback: inherit from master background
        if is_dark is None:
            is_dark = master_is_dark

        # Collect explicit title font size from layout placeholder (overrides txStyles default)
        title_size_pt: float | None = None
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0:
                for sz_tag in [_a("defRPr"), _a("rPr")]:
                    for el in ph.element.iter(sz_tag):
                        sz = el.get("sz")
                        if sz:
                            title_size_pt = _safe_pt(sz)
                            break
                    if title_size_pt:
                        break
                break

        layouts.append({
            "name":            layout.name,
            "background":      bg,
            "is_dark":         is_dark,
            "title_color":     title_color,
            "title_size_pt":   title_size_pt,   # explicit override; None = inherits txStyles
            "placeholder_count": len(ph_list),
            "placeholders":    ph_list,
        })
    return layouts


# ═════════════════════════════════════════════════════════════════════════════
# 6. MASTER NAMED SHAPES
# ═════════════════════════════════════════════════════════════════════════════

def extract_master_shapes(master,
                          scheme_resolver: dict[str, str]) -> list[dict]:
    """
    Extract every shape on the slide master with:
      name, shape_type, position/size in inches, fill descriptor,
      line color, and text content (for footer/label shapes).
    """
    shapes_out = []
    for shape in master.shapes:
        entry: dict = {
            "name":      shape.name,
            "shape_type": str(getattr(shape, "shape_type", "")).split(".")[-1],
            "left_in":  _safe_emu_in(getattr(shape, "left",  None)),
            "top_in":   _safe_emu_in(getattr(shape, "top",   None)),
            "width_in": _safe_emu_in(getattr(shape, "width", None)),
            "height_in":_safe_emu_in(getattr(shape, "height",None)),
        }

        # Fill — look inside spPr
        sp_pr = shape.element.find(_a("spPr"))
        if sp_pr is not None:
            fill = resolve_color(sp_pr, scheme_resolver)
            if fill:
                entry["fill"] = fill
            # Line color
            ln = sp_pr.find(_a("ln"))
            if ln is not None:
                lnf = resolve_color(ln, scheme_resolver)
                if lnf:
                    entry["line"] = lnf
                # Line width
                w = sp_pr.find(_a("ln"))
                if w is not None:
                    w_attr = w.get("w")
                    if w_attr:
                        entry["line_width_pt"] = round(int(w_attr) / 12700, 1)

        # Text content
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            try:
                text = shape.text_frame.text.strip()
                if text:
                    entry["text"] = text[:200]
            except Exception:
                pass

        shapes_out.append(entry)
    return shapes_out


# ═════════════════════════════════════════════════════════════════════════════
# 7. SLIDE SHAPE FILLS  (all slides)
# ═════════════════════════════════════════════════════════════════════════════

def extract_all_shape_fills(prs: Presentation,
                             scheme_resolver: dict[str, str]) -> dict:
    """
    Walk every shape on every slide, collect fill colors.
    Returns:
      {
        "hex_colors": Counter of hex → count (resolved where possible),
        "scheme_colors": Counter of scheme val → count,
        "large_bg_shapes": [{"slide": N, "name": ..., "fill": ..., "size_in": "WxH"}],
      }
    """
    hex_counter:    Counter = Counter()
    scheme_counter: Counter = Counter()
    large_shapes:   list    = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else ""
        for shape in slide.shapes:
            sp_pr = shape.element.find(_a("spPr"))
            if sp_pr is None:
                continue
            fill = resolve_color(sp_pr, scheme_resolver)
            if fill is None:
                continue
            if fill["type"] == "solid":
                if fill.get("hex"):
                    hex_counter[fill["hex"]] += 1
                if fill.get("scheme"):
                    scheme_counter[fill["scheme"]] += 1
            elif fill["type"] == "gradient":
                for stop in fill.get("stops", []):
                    if stop.get("hex"):
                        hex_counter[stop["hex"]] += 1
                    if stop.get("scheme"):
                        scheme_counter[stop["scheme"]] += 1

            # Flag large shapes (> half slide width or height) as potential backgrounds
            try:
                w = shape.width or 0
                h = shape.height or 0
                slide_w = prs.slide_width
                slide_h = prs.slide_height
                if w > slide_w * 0.5 or h > slide_h * 0.5:
                    large_shapes.append({
                        "slide": slide_idx,
                        "layout": layout_name,
                        "name": shape.name,
                        "fill": fill,
                        "width_in":  _safe_emu_in(w),
                        "height_in": _safe_emu_in(h),
                    })
            except Exception:
                pass

    return {
        "hex_colors":    dict(hex_counter.most_common(30)),
        "scheme_colors": dict(scheme_counter.most_common(20)),
        "large_bg_shapes": large_shapes,
    }


# ═════════════════════════════════════════════════════════════════════════════
# 8. TEXT STYLE CORPUS  (unchanged from v1, but scheme-color resolved)
# ═════════════════════════════════════════════════════════════════════════════

def extract_shapes_styles(shapes,
                          source_label: str,
                          scheme_resolver: dict[str, str],
                          font_resolver: dict[str, str]) -> list[dict]:
    """Recursively collect text-style records from a shape collection."""
    records: list[dict] = []
    for shape in shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                pf = para.font if hasattr(para, "font") else None
                for run in para.runs:
                    rf   = run.font
                    size = _safe_pt(str(int(rf.size.pt * 100)) if rf.size else None) or \
                           (_safe_pt(str(int(pf.size.pt * 100)) if pf and pf.size else None))
                    # color
                    def _font_color(font):
                        try:
                            if font and font.color and font.color.type is not None:
                                hex_c = f"#{font.color.rgb[0]:02X}{font.color.rgb[1]:02X}{font.color.rgb[2]:02X}"
                                return hex_c
                        except Exception:
                            pass
                        return None
                    color = _font_color(rf) or _font_color(pf)
                    # resolve font name
                    font_name = resolve_typeface(rf.name, font_resolver) or \
                                (resolve_typeface(pf.name, font_resolver) if pf else None)
                    records.append({
                        "text":       run.text[:80].strip(),
                        "font_name":  font_name,
                        "size_pt":    size,
                        "bold":       rf.bold,
                        "italic":     rf.italic,
                        "color_hex":  color,
                        "source":     source_label,
                    })
                if not para.runs and pf:
                    size = _safe_pt(str(int(pf.size.pt * 100)) if pf.size else None)
                    def _font_color2(font):
                        try:
                            if font and font.color and font.color.type is not None:
                                return f"#{font.color.rgb[0]:02X}{font.color.rgb[1]:02X}{font.color.rgb[2]:02X}"
                        except Exception:
                            pass
                        return None
                    color = _font_color2(pf)
                    font_name = resolve_typeface(pf.name, font_resolver) if pf else None
                    records.append({
                        "text":      para.text[:80].strip(),
                        "font_name": font_name,
                        "size_pt":   size,
                        "bold":      pf.bold,
                        "italic":    pf.italic,
                        "color_hex": color,
                        "source":    source_label,
                    })
        if hasattr(shape, "shapes"):
            records.extend(extract_shapes_styles(shape.shapes, source_label,
                                                  scheme_resolver, font_resolver))
    return records


# ═════════════════════════════════════════════════════════════════════════════
# 9. TYPOGRAPHY SUMMARY
# ═════════════════════════════════════════════════════════════════════════════

def build_typography_summary(all_records: list[dict],
                              tx_styles: dict,
                              font_theme: dict,
                              scheme_resolver: dict[str, str],
                              layouts: list[dict] | None = None) -> dict:
    """
    Derive canonical title/body/footer sizes and fonts.
    Prefer explicit values from tx_styles over heuristic corpus counts.
    """
    sizes  = Counter(r["size_pt"]   for r in all_records if r.get("size_pt"))
    fonts  = Counter(r["font_name"] for r in all_records if r.get("font_name"))
    colors = Counter(r["color_hex"] for r in all_records if r.get("color_hex"))

    # Prefer tx_styles values for title
    title_lvl1 = tx_styles.get("title", {}).get(1, {})
    body_lvl1  = tx_styles.get("body",  {}).get(1, {})

    title_rpr  = title_lvl1.get("default_run_props", {})
    body_rpr   = body_lvl1.get("default_run_props", {})

    title_size = title_rpr.get("size_pt") or \
                 (sorted(sizes.keys(), reverse=True)[0] if sizes else None)
    body_size  = body_rpr.get("size_pt") or \
                 (sorted(sizes.keys(), reverse=True)[1] if len(sizes) > 1 else None)

    # Footer: smallest size in corpus
    sorted_sizes = sorted(sizes.keys())
    footer_size = sorted_sizes[0] if sorted_sizes else None

    # Primary font: prefer theme minor font (body), then corpus
    primary_font = (font_theme.get("minor_latin") or
                    font_theme.get("major_latin") or
                    (fonts.most_common(1)[0][0] if fonts else None))

    # Title font: prefer theme major font
    title_font = font_theme.get("major_latin") or primary_font

    # Title color from tx_styles
    title_color_hex    = title_rpr.get("color_hex")
    title_color_scheme = title_rpr.get("color_scheme")

    # Body color from tx_styles
    body_color_hex    = body_rpr.get("color_hex")
    body_color_scheme = body_rpr.get("color_scheme")

    # Dominant fill colors — filter out noise before ranking:
    #   • pure black (#000000) and pure white (#FFFFFF) are text/bg defaults, not brand colors
    #   • near-black/near-white (luminance < 12 or > 243) are also noise
    #   • reds in the #C00000–#FF0000 band are typically CONFIDENTIAL labels / watermarks
    _NOISE_EXACT   = {"#000000", "#FFFFFF", "#000001", "#FFFFFE"}
    def _is_noise_color(h: str) -> bool:
        if h in _NOISE_EXACT:
            return True
        try:
            r = int(h[1:3], 16); g = int(h[3:5], 16); b = int(h[5:7], 16)
        except (ValueError, IndexError):
            return True
        lum = 0.299*r + 0.587*g + 0.114*b
        if lum < 12 or lum > 243:    # near-black or near-white
            return True
        # CONFIDENTIAL-style reds: high red, low green & blue
        if r > 160 and g < 40 and b < 40:
            return True
        return False
    dominant_colors = [c for c, _ in colors.most_common(30) if not _is_noise_color(c)][:15]

    # Collect any layout-level title size overrides (e.g. 40pt on dark title slides)
    layout_title_sizes: list[float] = []
    if layouts:
        for lay in layouts:
            sz = lay.get("title_size_pt")
            if sz:
                layout_title_sizes.append(sz)
    # The largest layout title size override (e.g. 40pt for full-bleed title layouts)
    title_size_max = max(layout_title_sizes) if layout_title_sizes else None

    return {
        "primary_font":        primary_font,
        "title_font":          title_font,
        "title_size_pt":       title_size,       # master txStyles default (e.g. 26pt)
        "title_size_max_pt":   title_size_max,   # largest layout override (e.g. 40pt for full-bleed)
        "title_cap":           title_rpr.get("cap"),
        "title_color_hex":     title_color_hex,
        "title_color_scheme":  title_color_scheme,
        "body_size_pt":        body_size,
        "body_font":           body_rpr.get("font") or primary_font,
        "body_color_hex":      body_color_hex,
        "body_color_scheme":   body_color_scheme,
        "footer_size_pt":      footer_size,
        "dominant_colors":     dominant_colors,
        "all_fonts":           dict(fonts.most_common(10)),
        "all_sizes_pt":        dict(sizes.most_common(20)),
        "all_colors":          dict(colors.most_common(20)),
    }


# ═════════════════════════════════════════════════════════════════════════════
# 10. SLIDE BACKGROUNDS (per slide)
# ═════════════════════════════════════════════════════════════════════════════

def extract_slide_backgrounds(prs: Presentation,
                               scheme_resolver: dict[str, str]) -> list[dict]:
    results = []
    for i, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else ""
        bg_el = slide.element.find(_p("bg"))
        bg = _parse_bg_element(bg_el, scheme_resolver)
        if bg["type"] == "inherited":
            # Report the layout name so agents know where to look
            bg["inherited_from"] = "layout"
        results.append({"slide": i, "layout": layout_name, "background": bg})
    return results


# ═════════════════════════════════════════════════════════════════════════════
# 11. TITLE BACKGROUND IMAGE EXTRACTION
# ═════════════════════════════════════════════════════════════════════════════

def extract_title_bg_jpeg(prs: Presentation, out_path: Path) -> bool:
    """Extract the picture placeholder image from the title layout."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if "title" in layout.name.lower() or "full picture" in layout.name.lower():
                for ph in layout.placeholders:
                    try:
                        image = ph.image
                        if image:
                            out_path.write_bytes(image.blob)
                            return True
                    except Exception:
                        pass
    # fallback: first image shape on first slide
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "image"):
                try:
                    out_path.write_bytes(shape.image.blob)
                    return True
                except Exception:
                    pass
    return False


# ═════════════════════════════════════════════════════════════════════════════
# 12. FULL EXTRACTION ORCHESTRATOR
# ═════════════════════════════════════════════════════════════════════════════

def extract_full_style(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))

    slide_w = _safe_emu_in(prs.slide_width)
    slide_h = _safe_emu_in(prs.slide_height)
    aspect  = f"{round(slide_w / slide_h * 9)}:9" if slide_h else "16:9"

    master = prs.slide_masters[0]

    # ── Theme ────────────────────────────────────────────────────────────────
    theme_colors  = extract_theme_colors(master)
    font_theme    = extract_theme_fonts(master)
    scheme_resolver = build_scheme_resolver(theme_colors)
    font_resolver   = font_theme.get("resolved", {})

    # ── Master level ─────────────────────────────────────────────────────────
    tx_styles       = extract_tx_styles(master, scheme_resolver, font_resolver)
    master_bg       = extract_master_background(master, scheme_resolver)
    master_shapes   = extract_master_shapes(master, scheme_resolver)

    # Determine master darkness so layouts that fully inherit can default to it
    master_is_dark  = _is_dark_fill(master_bg, scheme_resolver)
    # bg_ref idx=1001 with scheme lt1/bg1 = light; scheme dk1/tx1 = dark
    if master_is_dark is None and master_bg.get("type") == "bg_ref":
        scheme = master_bg.get("scheme", "")
        if scheme in ("lt1", "lt2", "bg1", "bg2"):
            master_is_dark = False
        elif scheme in ("dk1", "dk2", "tx1", "tx2"):
            master_is_dark = True
    # Default: if master bg is fully unresolved, assume light (most common)
    if master_is_dark is None:
        master_is_dark = False

    # ── Layout level ─────────────────────────────────────────────────────────
    layouts         = extract_layout_info(master, scheme_resolver, font_resolver, master_is_dark)

    # ── Slide level ──────────────────────────────────────────────────────────
    slide_bgs       = extract_slide_backgrounds(prs, scheme_resolver)
    shape_fills     = extract_all_shape_fills(prs, scheme_resolver)

    # ── Text corpus ──────────────────────────────────────────────────────────
    all_records: list = []
    for slide in prs.slides:
        all_records.extend(extract_shapes_styles(slide.shapes,   "slide",  scheme_resolver, font_resolver))
    master_records: list = extract_shapes_styles(master.shapes,  "master", scheme_resolver, font_resolver)
    layout_records: list = []
    for layout in master.slide_layouts:
        layout_records.extend(extract_shapes_styles(layout.shapes, "layout", scheme_resolver, font_resolver))

    combined = all_records + master_records + layout_records

    # ── Typography summary ───────────────────────────────────────────────────
    typography = build_typography_summary(combined, tx_styles, font_theme, scheme_resolver, layouts)

    return {
        "slide_dimensions": {
            "width_in":   slide_w,
            "height_in":  slide_h,
            "aspect_ratio": aspect,
            "width_emu":  int(prs.slide_width),
            "height_emu": int(prs.slide_height),
        },
        "theme": {
            "color_scheme": theme_colors,
            "font_scheme":  font_theme,
            "scheme_name":  list(theme_colors.values())[0].get("scheme_name", "") if theme_colors else "",
        },
        "tx_styles":       tx_styles,
        "master_background": master_bg,
        "master_shapes":   master_shapes,
        "layouts":         layouts,
        "slide_backgrounds": slide_bgs,
        "shape_fills":     shape_fills,
        "typography":      typography,
        "sample_text_styles": [r for r in combined if r.get("text")][:40],
        "extraction_stats": {
            "slide_shape_records":  len(all_records),
            "master_shape_records": len(master_records),
            "layout_shape_records": len(layout_records),
            "total_slides":         len(prs.slides),
            "total_layouts":        len(layouts),
            "theme_slots_found":    len(theme_colors),
            "font_scheme_found":    bool(font_theme),
            "tx_styles_found":      bool(tx_styles.get("title") or tx_styles.get("body")),
        },
    }


# ═════════════════════════════════════════════════════════════════════════════
# 13. BRAND GUIDE GENERATION
# ═════════════════════════════════════════════════════════════════════════════

def generate_brand_guide_md(style: dict, name: str, source_file: str,
                            is_builtin: bool = False) -> str:
    t     = style["typography"]
    dims  = style["slide_dimensions"]
    theme = style.get("theme", {})
    clr   = theme.get("color_scheme", {})
    fonts = theme.get("font_scheme", {})
    tx    = style.get("tx_styles", {})
    mshapes = style.get("master_shapes", [])
    layouts = style.get("layouts", [])
    now   = datetime.now().strftime("%Y-%m-%d")

    primary_color = _pick_primary_color(style, is_builtin_bmw_ci=is_builtin)
    major_font    = fonts.get("major_latin", "Arial")
    minor_font    = fonts.get("minor_latin", "Arial")
    dark_layouts  = [l for l in layouts if l.get("is_dark") is True]
    light_layouts = [l for l in layouts if l.get("is_dark") is False]

    lines = [
        f"# Brand Styling Guide — {name}",
        f"",
        f"Extracted from `{Path(source_file).name}` on {now}.",
        f"",
        f"---",
        f"",
    ]

    # ── Style overview (generic, not brand-specific) ──────────────────────────
    lines += [
        f"## Style Overview",
        f"",
        f"Extracted from `{Path(source_file).name}` on {datetime.now().strftime('%Y-%m-%d')}.",
        f"",
        f"### Primary Brand Color",
        f"",
        f"- **Primary color:** `{primary_color}` — used for titles, headers, and key accents.",
        f"- Supporting palette: " + " / ".join(f"`{c}`" for c in [c for c in t.get("dominant_colors", []) if c != primary_color][:5]),
        f"",
        f"### Font Family",
        f"",
        f"| Role | Font |",
        f"|------|------|",
        f"| Headings (major) | `{major_font}` |",
        f"| Body (minor) | `{minor_font}` |",
        f"",
        f"### Layout Summary",
        f"",
        f"- Total layouts: {len(layouts)}",
        f"- Dark layouts: {len(dark_layouts)}",
        f"- Light layouts: {len(light_layouts)}",
        f"",
        f"---",
        f"",
    ]

    # ── Slide dimensions ─────────────────────────────────────────────────────
    lines += [
        f"## Slide Dimensions",
        f"",
        f"- **Width:** {dims['width_in']}\" ({dims['aspect_ratio']})",
        f"- **Height:** {dims['height_in']}\"",
        f"- **Width EMU:** {dims['width_emu']}",
        f"- **Height EMU:** {dims['height_emu']}",
        f"",
        f"---",
        f"",
    ]

    # ── Theme color scheme ───────────────────────────────────────────────────
    lines += [
        f"## Theme Color Scheme — `{theme.get('scheme_name', 'Unknown')}`",
        f"",
        f"These are the canonical slot values from `<a:clrScheme>`. Every shape that uses a",
        f"scheme color (e.g. `tx2`, `accent1`) inherits from this table.",
        f"",
        f"| Slot | Hex | Label |",
        f"|------|-----|-------|",
    ]
    alias_note = {"dk1": " ← also `tx1`", "lt1": " ← also `bg1`",
                  "dk2": " ← also `tx2`", "lt2": " ← also `bg2`"}
    for slot in CLR_SCHEME_SLOTS:
        info = clr.get(slot, {})
        hex_c = info.get("hex", "—")
        label = CLR_LABELS.get(slot, slot) + alias_note.get(slot, "")
        lines.append(f"| `{slot}` | `{hex_c}` | {label} |")

    lines += [
        f"",
        f"---",
        f"",
    ]

    # ── Font scheme ──────────────────────────────────────────────────────────
    lines += [
        f"## Theme Font Scheme — `{fonts.get('scheme_name', 'Unknown')}`",
        f"",
        f"| Token | Resolved Font | Role |",
        f"|-------|---------------|------|",
        f"| `+mj-lt` | `{fonts.get('major_latin', '—')}` | Headings / titles |",
        f"| `+mn-lt` | `{fonts.get('minor_latin', '—')}` | Body text |",
    ]
    if fonts.get("major_ea"):
        lines.append(f"| `+mj-ea` | `{fonts['major_ea']}` | East-Asian headings |")
    if fonts.get("minor_ea"):
        lines.append(f"| `+mn-ea` | `{fonts['minor_ea']}` | East-Asian body |")
    lines += [f"", f"---", f""]

    # ── Typography from txStyles ─────────────────────────────────────────────
    lines += [
        f"## Master Text Styles (from `p:txStyles`)",
        f"",
        f"### Title Style",
        f"",
    ]
    title_lvl = tx.get("title", {})
    if title_lvl:
        lines += [
            f"| Attr | Value |",
            f"|------|-------|",
        ]
        lvl1 = title_lvl.get(1, {})
        rpr  = lvl1.get("default_run_props", {})
        if rpr.get("size_pt"):  lines.append(f"| Font size | **{rpr['size_pt']}pt** |")
        if rpr.get("font"):     lines.append(f"| Font | `{rpr['font']}` |")
        if rpr.get("cap"):      lines.append(f"| Capitalization | `{rpr['cap']}` |")
        if rpr.get("bold") is not None: lines.append(f"| Bold | {rpr['bold']} |")
        color_str = rpr.get('color_hex') or (f"scheme:`{rpr['color_scheme']}`" if rpr.get('color_scheme') else "—")
        lines.append(f"| Color | `{color_str}` |")
        if lvl1.get("align"):   lines.append(f"| Alignment | `{lvl1['align']}` |")
        bu = lvl1.get("bullet", {})
        if bu.get("type") == "none": lines.append(f"| Bullet | none |")
        lines += [f""]
    else:
        lines += [f"*Not found in master (may be in theme defaults).*", f""]

    lines += [f"### Body Style — Bullet Levels", f""]
    body_lvls = tx.get("body", {})
    if body_lvls:
        lines += [
            f"| Lvl | Size | Font | Bullet | Bullet Font | Bullet Color | Space After | Indent (EMU) |",
            f"|-----|------|------|--------|-------------|--------------|-------------|--------------|",
        ]
        for lvl_num in sorted(body_lvls.keys()):
            lvl   = body_lvls[lvl_num]
            rpr   = lvl.get("default_run_props", {})
            bu    = lvl.get("bullet", {})
            size  = rpr.get("size_pt", "—")
            font  = rpr.get("font", "—")
            bu_char  = bu.get("char", "none") if bu.get("type") == "char" else bu.get("type", "—")
            bu_font  = bu.get("font", "—") if bu.get("type") == "char" else "—"
            bu_color = bu.get("color_hex") or (f"scheme:{bu.get('color_scheme')}" if bu.get("color_scheme") else "—")
            spc_aft  = lvl.get("space_after_pt", "—")
            indent   = lvl.get("margin_left_emu", "—")
            lines.append(f"| {lvl_num} | {size} | `{font}` | `{bu_char}` | `{bu_font}` | `{bu_color}` | {spc_aft} | {indent} |")
        lines += [f""]
    else:
        lines += [f"*Not found in master.*", f""]

    lines += [f"---", f""]

    # ── Master background ────────────────────────────────────────────────────
    mb = style.get("master_background", {})
    lines += [
        f"## Master Background",
        f"",
        f"- **Type:** `{mb.get('type', 'unknown')}`",
    ]
    if mb.get("hex"):    lines.append(f"- **Color:** `{mb['hex']}`")
    if mb.get("scheme"): lines.append(f"- **Scheme slot:** `{mb['scheme']}`")
    if mb.get("idx"):    lines.append(f"- **bgRef idx:** `{mb['idx']}`")
    if mb.get("stops"):
        stops_str = ", ".join(f"`{s.get('hex') or s.get('scheme')}`" for s in mb["stops"][:3])
        lines.append(f"- **Gradient stops:** {stops_str}")
    lines += [f"", f"---", f""]

    # ── Master named shapes ──────────────────────────────────────────────────
    lines += [
        f"## Master Shapes (footer, lines, logos)",
        f"",
        f"| Name | Type | Position | Size | Fill | Text |",
        f"|------|------|----------|------|------|------|",
    ]
    for sh in mshapes:
        pos  = f"{sh.get('left_in','?')}\"×{sh.get('top_in','?')}\""
        size = f"{sh.get('width_in','?')}\"×{sh.get('height_in','?')}\""
        fill_desc = "—"
        f = sh.get("fill")
        if f:
            if f.get("type") == "solid":
                fill_desc = f.get("hex") or f"scheme:{f.get('scheme')}"
            else:
                fill_desc = f.get("type", "—")
        text = (sh.get("text") or "")[:40]
        lines.append(
            f"| `{sh['name']}` | {sh.get('shape_type','?')} | {pos} | {size} | `{fill_desc}` | {text} |"
        )
    lines += [f"", f"---", f""]

    # ── Layouts ──────────────────────────────────────────────────────────────
    lines += [
        f"## Layouts",
        f"",
        f"| Layout Name | Background | Dark? | Title Color | Title Size | Placeholders |",
        f"|-------------|-----------|-------|-------------|------------|--------------|",
    ]
    for lay in layouts:
        bg = lay.get("background", {})
        if bg.get("type") == "solid":
            bg_desc = bg.get("hex") or f"scheme:{bg.get('scheme')}"
        elif bg.get("type") == "gradient":
            stops = bg.get("stops", [])
            bg_desc = "gradient(" + ", ".join(
                s.get("hex") or f"scheme:{s.get('scheme')}" for s in stops[:2]
            ) + ")"
        elif bg.get("type") == "picture":
            bg_desc = "picture/image"
        elif bg.get("type") == "bg_ref":
            bg_desc = f"bgRef(scheme:{bg.get('scheme','')})"
        else:
            bg_desc = "inherited"

        dark_str = "yes" if lay.get("is_dark") is True else \
                   "no"  if lay.get("is_dark") is False else "?"
        tc = lay.get("title_color") or {}
        tc_str = tc.get("hex") or (f"scheme:{tc.get('scheme')}" if tc.get("scheme") else "inherited")
        ts_str = f"{lay.get('title_size_pt')}pt" if lay.get("title_size_pt") else "inherited"
        lines.append(
            f"| `{lay['name']}` | {bg_desc} | {dark_str} | {tc_str} | {ts_str} | {lay['placeholder_count']} |"
        )
    lines += [f"", f"---", f""]

    # ── Shape fill palette ───────────────────────────────────────────────────
    sf = style.get("shape_fills", {})
    hex_fills = sf.get("hex_colors", {})
    scheme_fills = sf.get("scheme_colors", {})
    lines += [
        f"## Shape Fill Colors (across all slides)",
        f"",
        f"These are the actual fill colors used on slide shapes — the true brand palette.",
        f"",
        f"### Hardcoded Hex Colors (by frequency)",
        f"",
        f"| Hex | Count |",
        f"|-----|-------|",
    ]
    for hex_c, count in list(hex_fills.items())[:15]:
        lines.append(f"| `{hex_c}` | {count} |")

    lines += [
        f"",
        f"### Scheme-color Fills (by frequency)",
        f"",
        f"| Scheme Slot | Count | Resolved Hex |",
        f"|------------|-------|--------------|",
    ]
    resolver = build_scheme_resolver(clr)
    for sc, count in list(scheme_fills.items())[:10]:
        resolved = resolver.get(sc, "—")
        lines.append(f"| `{sc}` | {count} | `{resolved}` |")

    large = sf.get("large_bg_shapes", [])
    if large:
        lines += [
            f"",
            f"### Large Background Shapes",
            f"",
            f"| Slide | Layout | Shape Name | Fill |",
            f"|-------|--------|------------|------|",
        ]
        for s in large[:10]:
            f_desc = "—"
            f = s.get("fill", {})
            if f.get("type") == "solid":
                f_desc = f.get("hex") or f"scheme:{f.get('scheme')}"
            elif f.get("type") == "gradient":
                f_desc = "gradient"
            lines.append(f"| {s['slide']} | `{s['layout']}` | `{s['name']}` | `{f_desc}` |")

    lines += [f"", f"---", f""]

    # ── Typography summary ───────────────────────────────────────────────────
    lines += [
        f"## Typography Summary",
        f"",
        f"| Element | Font | Size | Notes | Color |",
        f"|---------|------|------|-------|-------|",
        f"| **Title (default)** | `{t.get('title_font') or '—'}` | {t.get('title_size_pt') or '—'}pt"
        f" | master txStyles default | `{t.get('title_color_hex') or t.get('title_color_scheme') or '—'}` |",
        *(
            [f"| **Title (full-bleed)** | `{t.get('title_font') or '—'}` | {t.get('title_size_max_pt')}pt"
             f" | layout override (dark/picture title slides) | `{t.get('title_color_hex') or t.get('title_color_scheme') or '—'}` |"]
            if t.get("title_size_max_pt") and t.get("title_size_max_pt") != t.get("title_size_pt") else []
        ),
        f"| **Body** | `{t.get('body_font') or t.get('primary_font') or '—'}` | {t.get('body_size_pt') or '—'}pt"
        f" | | `{t.get('body_color_hex') or t.get('body_color_scheme') or '—'}` |",
        f"| **Footer** | `{t.get('primary_font') or '—'}` | {t.get('footer_size_pt') or '—'}pt | | — |",
        f"",
        f"---",
        f"",
        f"## Usage Notes",
        f"",
    ]
    # Resolve primary brand color: dk2 > large bg shape fill > accent1 > corpus
    _clr_scheme   = theme.get("color_scheme", {})
    _primary_brand = (_clr_scheme.get("dk2",     {}).get("hex") or
                      _clr_scheme.get("accent1", {}).get("hex") or
                      (list(hex_fills.keys())[0] if hex_fills else None) or
                      (t.get("dominant_colors") or ["—"])[0])
    lines += [
        f"1. **Primary brand color** (theme `dk2`): `{_primary_brand}`",
        f"2. **Title font:** `{t.get('title_font') or '—'}` — resolves from theme token `+mj-lt`",
        f"3. **Body font:** `{t.get('body_font') or t.get('primary_font') or '—'}` — resolves from theme token `+mn-lt`",
        f"4. **Bullet (level 1):** char `{body_lvls.get(1, {}).get('bullet', {}).get('char', '?') if body_lvls else '?'}`"
        f" in `{body_lvls.get(1, {}).get('bullet', {}).get('font', '?') if body_lvls else '?'}`"
        f" — color `{body_lvls.get(1, {}).get('bullet', {}).get('color_hex', '?') if body_lvls else '?'}`",
        f"5. **Title scheme color:** `{t.get('title_color_scheme', '—')}` → `{resolver.get(t.get('title_color_scheme',''), '—')}`",
        f"6. See `style.json` for the complete machine-readable profile.",
    ]

    return "\n".join(lines)


# ═════════════════════════════════════════════════════════════════════════════
# REGISTRY OPERATIONS  (unchanged from v1)
# ═════════════════════════════════════════════════════════════════════════════

def slugify(text: str) -> str:
    text = text.lower().strip()
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"[\s_-]+", "-", text)
    return text[:64]


def load_registry() -> dict:
    if REGISTRY_FILE.exists():
        try:
            return json.loads(REGISTRY_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"version": 1, "styles": {}}


def save_registry(registry: dict) -> None:
    REGISTRY_FILE.write_text(
        json.dumps(registry, indent=2, ensure_ascii=False), encoding="utf-8"
    )


def file_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def bootstrap_bmw_ci() -> None:
    """No-op — Dela doesn't ship a built-in BMW CI style.
    Styles are added by cloning .pptx files via clone_pptx_style tool.
    Kept for API compatibility with the original script.
    """
    pass


def _is_noise_hex(h: str) -> bool:
    """Return True if this hex color is pure black/white, near-black/white, or CONFIDENTIAL red."""
    _NOISE_EXACT = {"#000000", "#FFFFFF"}
    if not h or not h.startswith("#") or len(h) != 7:
        return True
    if h in _NOISE_EXACT:
        return True
    try:
        r = int(h[1:3], 16); g = int(h[3:5], 16); b = int(h[5:7], 16)
    except ValueError:
        return True
    lum = 0.299*r + 0.587*g + 0.114*b
    if lum < 12 or lum > 243:
        return True
    if r > 160 and g < 40 and b < 40:   # CONFIDENTIAL-red band
        return True
    return False


def _pick_primary_color(style: dict, is_builtin_bmw_ci: bool = False) -> str | None:
    """
    Choose the most meaningful primary brand color for menu display and brand guide.

    BMW CI (built-in):  always returns Ocean Blue #035970 (dk2).

    Cloned styles — priority:
      1. Most-common non-noise background fill on *dark* layouts (the deck's visual identity)
      2. Large background shape fills (solid, non-noise)
      3. Title color (from txStyles) — if non-noise and non-white
      4. Most-common non-noise hex from dominant_colors corpus
      5. theme dk2 (last resort fallback)
    """
    clr = style.get("theme", {}).get("color_scheme", {})
    dk2 = clr.get("dk2", {}).get("hex")

    # 1. Dark layout backgrounds
    for lay in style.get("layouts", []):
        if lay.get("is_dark") is True:
            bg = lay.get("background", {})
            if bg.get("type") == "solid":
                h = bg.get("hex")
                if h and not _is_noise_hex(h):
                    return h

    # 2. Large background shapes
    large = style.get("shape_fills", {}).get("large_bg_shapes", [])
    for s in large:
        f = s.get("fill", {})
        if f.get("type") == "solid":
            h = f.get("hex")
            if h and not _is_noise_hex(h):
                return h

    # 3. Title color from txStyles (the color titles are actually rendered in)
    title_hex = style.get("typography", {}).get("title_color_hex")
    if title_hex and not _is_noise_hex(title_hex):
        return title_hex

    # 4. First non-noise color from dominant corpus
    for c in style.get("typography", {}).get("dominant_colors", []):
        if not _is_noise_hex(c):
            return c

    # 5. dk2 fallback
    return dk2


# Keep the old name as an alias for backwards compat in register_style
def _pick_dominant_color(style: dict) -> str | None:
    return _pick_primary_color(style, is_builtin_bmw_ci=False)


def register_style(slug: str, name: str, description: str,
                   source_path: Path, style_dir: Path, style: dict) -> None:
    registry = load_registry()
    t = style["typography"]
    is_builtin = False  # Dela doesn't ship built-in styles
    primary_color = _pick_primary_color(style, is_builtin_bmw_ci=False)

    # Layout statistics
    layouts = style.get("layouts", [])
    dark_count  = sum(1 for l in layouts if l.get("is_dark") is True)
    light_count = sum(1 for l in layouts if l.get("is_dark") is False)

    # Font family name (e.g. "BMWGroupTN Condensed" → "BMWGroupTN")
    major_font = style["theme"]["font_scheme"].get("major_latin", "")
    minor_font = style["theme"]["font_scheme"].get("minor_latin", "")

    registry["styles"][slug] = {
        "slug":         slug,
        "name":         name,
        "description":  description,
        "is_builtin":   False,              # All styles are user-cloned
        "source_file":  source_path.name,
        "source_sha256": file_sha256(source_path),
        "cloned_at":    datetime.now(timezone.utc).isoformat(),
        "style_dir":    str(style_dir),
        "brand_guide_path": str(style_dir / "brand-guide.md"),
        "has_title_bg":  (style_dir / "title-bg.jpeg").exists(),
        "slide_dimensions": style["slide_dimensions"],
        "typography_summary": {
            "primary_font":    t.get("primary_font"),
            "title_font":      major_font,
            "title_size_pt":   t.get("title_size_pt"),
            "body_size_pt":    t.get("body_size_pt"),
            "primary_color":   primary_color,   # meaningful brand color for this style
            "dominant_color":  primary_color,   # kept for backward compat
            "palette":         t.get("dominant_colors", [])[:6],
        },
        "theme_summary": {
            "scheme_name": style["theme"].get("scheme_name", ""),
            "dk1":  style["theme"]["color_scheme"].get("dk1",  {}).get("hex"),
            "lt1":  style["theme"]["color_scheme"].get("lt1",  {}).get("hex"),
            "dk2":  style["theme"]["color_scheme"].get("dk2",  {}).get("hex"),
            "accent1": style["theme"]["color_scheme"].get("accent1", {}).get("hex"),
            "major_font": major_font,
            "minor_font": minor_font,
        },
        "layout_stats": {
            "total":  len(layouts),
            "dark":   dark_count,
            "light":  light_count,
        },
        "total_slides_in_source": style["extraction_stats"]["total_slides"],
        "layout_count":           style["extraction_stats"]["total_layouts"],
    }
    save_registry(registry)


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Clone and persist the full visual style of a .pptx file into the style registry."
    )
    parser.add_argument("pptx",          help="Path to source .pptx file")
    parser.add_argument("--name",        default="", help="Human-readable style name")
    parser.add_argument("--description", default="", help="Short description")
    parser.add_argument("--overwrite",   action="store_true",
                        help="Overwrite existing style with same slug")
    args = parser.parse_args()

    source = Path(args.pptx).expanduser().resolve()
    if not source.exists():
        print(f"ERROR: File not found: {source}"); sys.exit(1)
    if source.suffix.lower() != ".pptx":
        print(f"ERROR: Expected .pptx, got: {source.name}"); sys.exit(1)

    name        = args.name.strip() or source.stem
    slug        = slugify(name)
    description = args.description.strip() or f"Style cloned from {source.name}"

    print(f"[clone] Source:      {source}")
    print(f"[clone] Style name:  {name}")
    print(f"[clone] Slug:        {slug}")

    REGISTRY_DIR.mkdir(parents=True, exist_ok=True)
    bootstrap_bmw_ci()

    registry = load_registry()
    if slug in registry["styles"] and not args.overwrite:
        if registry["styles"][slug].get("is_builtin"):
            print(f"ERROR: '{slug}' is a built-in style and cannot be overwritten.")
            sys.exit(1)
        print(f"WARNING: Style '{slug}' already exists. Use --overwrite to replace it.")
        sys.exit(1)

    style_dir = REGISTRY_DIR / slug
    style_dir.mkdir(parents=True, exist_ok=True)

    print("[clone] Extracting style profile (full master + theme + layout pass)...")
    style = extract_full_style(source)

    stats = style["extraction_stats"]
    print(f"[clone]   theme slots:    {stats['theme_slots_found']}/12")
    print(f"[clone]   font scheme:    {'yes' if stats['font_scheme_found'] else 'no'}")
    print(f"[clone]   txStyles:       {'yes' if stats['tx_styles_found'] else 'no'}")
    print(f"[clone]   layouts:        {stats['total_layouts']}")
    print(f"[clone]   slides:         {stats['total_slides']}")

    style_json_path = style_dir / "style.json"
    style_json_path.write_text(json.dumps(style, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"[clone] style.json   → {style_json_path}")

    is_builtin_clone = False  # Dela doesn't ship built-in styles
    brand_guide_md  = generate_brand_guide_md(style, name, str(source), is_builtin=is_builtin_clone)
    brand_guide_path = style_dir / "brand-guide.md"
    brand_guide_path.write_text(brand_guide_md, encoding="utf-8")
    print(f"[clone] brand-guide  → {brand_guide_path}")

    print("[clone] Attempting title background extraction...")
    title_bg_path = style_dir / "title-bg.jpeg"
    if extract_title_bg_jpeg(Presentation(str(source)), title_bg_path):
        print(f"[clone] title-bg     → {title_bg_path}")
    else:
        print("[clone] title-bg     — not found")

    source_copy = style_dir / "source.pptx"
    if source.resolve() != source_copy.resolve():
        shutil.copy2(source, source_copy)
        print(f"[clone] source copy  → {source_copy}")
    else:
        print(f"[clone] source copy  — skipped (source is already registry copy)")

    register_style(slug, name, description, source, style_dir, style)
    print(f"[clone] registry     → {REGISTRY_FILE}")

    t = style["typography"]
    sf = style.get("shape_fills", {})
    shape_hex = list(sf.get("hex_colors", {}).keys())
    print(f"\n✓ Style '{name}' cloned successfully (slug: {slug})")
    print(f"  Directory:      {style_dir}")
    print(f"  Theme scheme:   {style['theme'].get('scheme_name', '—')}")
    print(f"  Title font:     {t.get('title_font') or '—'}  {t.get('title_size_pt') or '—'}pt")
    print(f"  Body font:      {t.get('body_font') or t.get('primary_font') or '—'}  {t.get('body_size_pt') or '—'}pt")
    print(f"  Primary brand:  {_pick_dominant_color(style) or '—'}  (theme dk2)")
    print(f"\n  Use it: tell the presentation builder \"use the '{name}' style\"")


if __name__ == "__main__":
    main()
