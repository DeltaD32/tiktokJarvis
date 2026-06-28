"""
pptx_lib.infographics — Timeline, RAG dashboard, pull-quote, progress ring,
hexagon grid, Venn diagram, funnel, callout, hero number.
"""

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Pt

from .constants import (
    COLOR_ACCENT,
    COLOR_PRIMARY,
    COLOR_SECONDARY,
    FREE_HEIGHT,
    FREE_LEFT,
    FREE_TOP,
    FREE_WIDTH,
    PALETTE,
    WHITE,
)


def add_hero_number(
    slide,
    number,
    subtitle="",
    left=None,
    top=None,
    width=None,
    height=None,
    number_size=Pt(48),
    subtitle_size=Pt(24),
    number_color=None,
    subtitle_color=None,
):
    """Large centered number with optional subtitle below.

    Defaults to full FREE area if position not specified.
    """
    left = left or FREE_LEFT
    top = top or FREE_TOP
    width = width or FREE_WIDTH
    height = height or FREE_HEIGHT
    number_color = number_color or COLOR_PRIMARY
    subtitle_color = subtitle_color or COLOR_SECONDARY

    txBox = slide.shapes.add_textbox(left, top, width, Emu(int(height * 0.6)))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = number_size
    p.font.bold = True
    p.font.color.rgb = number_color
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            left, top + Emu(int(height * 0.6)), width, Emu(int(height * 0.4))
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = subtitle_size
        p2.font.color.rgb = subtitle_color
        p2.alignment = PP_ALIGN.CENTER


def add_timeline(
    slide,
    milestones,
    left,
    top,
    width,
    line_color=COLOR_PRIMARY,
    dot_color=COLOR_ACCENT,
):
    """Horizontal milestone timeline with labels alternating above/below.

    milestones: list of (label, date_str) tuples in chronological order.
    """
    line_y = top + Cm(2.5)
    dot_r = Cm(0.22)
    lbl_h = Cm(1.0)
    lbl_w = Cm(3.0)

    n = len(milestones)
    spacing = Emu(int(int(width) / (n + 1)))

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, line_y, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = line_color
    line.line.fill.background()

    for i, (label, date_str) in enumerate(milestones):
        x = left + Emu((i + 1) * int(spacing))
        above = i % 2 == 0

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x - dot_r, line_y - dot_r, 2 * dot_r, 2 * dot_r
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.color.rgb = line_color
        dot.line.width = Pt(1)

        tick_top = line_y - Cm(1.0) if above else line_y
        tick_h = Cm(1.0)
        tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Pt(0.5), tick_top, Pt(1), tick_h)
        tick.fill.solid()
        tick.fill.fore_color.rgb = line_color
        tick.line.fill.background()

        l_top = (line_y - Cm(1.0) - lbl_h) if above else (line_y + Cm(1.0))
        l_left = x - Emu(int(lbl_w) // 2)
        lbl = slide.shapes.add_textbox(l_left, l_top, lbl_w, lbl_h)
        tf = lbl.text_frame
        tf.word_wrap = True
        pl = tf.paragraphs[0]
        pl.text = label
        pl.font.size = Pt(11)
        pl.font.bold = True
        pl.font.color.rgb = COLOR_PRIMARY
        pl.alignment = PP_ALIGN.CENTER

        date_top = l_top + lbl_h if above else l_top - Cm(0.5)
        date_box = slide.shapes.add_textbox(l_left, date_top, lbl_w, Cm(0.5))
        tf2 = date_box.text_frame
        pd = tf2.paragraphs[0]
        pd.text = date_str
        pd.font.size = Pt(9)
        pd.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
        pd.alignment = PP_ALIGN.CENTER


def add_swimlane_timeline(
    slide, swimlanes, left, top, width, row_height=Cm(2.2), line_color=COLOR_PRIMARY
):
    """Multi-row swimlane timeline, one lane per workstream.

    swimlanes: list of (lane_name, [(label, x_fraction), ...]) tuples.
    x_fraction: 0.0 = leftmost, 1.0 = rightmost.
    """
    label_col_w = Cm(3.0)
    timeline_w = Emu(int(width) - int(label_col_w))
    dot_r = Cm(0.2)

    for row_i, (lane_name, events) in enumerate(swimlanes):
        row_top = top + Emu(row_i * int(row_height))
        line_y = row_top + Emu(int(row_height) // 2)
        tl_left = left + label_col_w

        lbl = slide.shapes.add_textbox(left, row_top, label_col_w, row_height)
        tf = lbl.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pl = tf.paragraphs[0]
        pl.text = lane_name
        pl.font.size = Pt(11)
        pl.font.bold = True
        pl.font.color.rgb = COLOR_PRIMARY

        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tl_left, line_y, timeline_w, Pt(1.5))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        ln.line.fill.background()

        for label, x_frac in events:
            x = tl_left + Emu(int(int(timeline_w) * x_frac))
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x - dot_r, line_y - dot_r, 2 * dot_r, 2 * dot_r
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = line_color
            dot.line.fill.background()

            ev_lbl = slide.shapes.add_textbox(x - Cm(1.5), line_y - Cm(1.2), Cm(3.0), Cm(0.9))
            tf2 = ev_lbl.text_frame
            tf2.word_wrap = True
            pe = tf2.paragraphs[0]
            pe.text = label
            pe.font.size = Pt(9)
            pe.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            pe.alignment = PP_ALIGN.CENTER


RAG_COLORS = {
    "R": RGBColor(0xC0, 0x39, 0x2B),
    "A": RGBColor(0xE6, 0x7E, 0x22),
    "G": RGBColor(0x1A, 0x7A, 0x3C),
}


def add_rag_dashboard(slide, items, left, top, width, row_height=Cm(1.25), dot_diameter=Cm(0.65)):
    """RAG status dashboard with dot, topic, owner, and note per row.

    items: list of (topic, status, owner, note) tuples.
    status: "R", "A", or "G".
    """
    from .helpers import send_to_back

    dot_d = dot_diameter
    dot_pad = Cm(0.3)
    col_topic_w = Cm(5.0)
    col_owner_w = Cm(3.5)
    col_note_w = Emu(
        int(width) - int(dot_d) - int(dot_pad * 3) - int(col_topic_w) - int(col_owner_w)
    )

    for i, (topic, status, owner, note) in enumerate(items):
        row_top = top + Emu(i * int(row_height))
        row_mid = row_top + Emu(int(row_height) // 2)

        if i % 2 == 0:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, row_top, width, row_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFA)
            bg.line.fill.background()
            send_to_back(slide, bg)

        dot_top = row_mid - Emu(int(dot_d) // 2)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + dot_pad, dot_top, dot_d, dot_d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = RAG_COLORS.get(status, RAG_COLORS["A"])
        dot.line.fill.background()

        cursor = left + dot_pad + dot_d + dot_pad

        tb_topic = slide.shapes.add_textbox(cursor, row_top, col_topic_w, row_height)
        tf = tb_topic.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        cursor += col_topic_w

        tb_owner = slide.shapes.add_textbox(cursor, row_top, col_owner_w, row_height)
        tf2 = tb_owner.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.text = owner
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        cursor += col_owner_w

        tb_note = slide.shapes.add_textbox(cursor, row_top, col_note_w, row_height)
        tf3 = tb_note.text_frame
        tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_rag_header(slide, left, top, width, row_height=Cm(1.0)):
    """Dark header row above a RAG dashboard."""
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, row_height)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = COLOR_PRIMARY
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "  ●   Topic                          Owner        Status / Notes"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return hdr


def add_pull_quote(
    slide,
    quote,
    attribution,
    left,
    top,
    width,
    height,
    accent_color=COLOR_ACCENT,
    text_color=COLOR_PRIMARY,
):
    """Large centered quote with accent bar and attribution.

    quote: str — keep to 2-3 lines max.
    attribution: str — e.g. "— Satya Nadella, Microsoft CEO".
    """
    bar_w = Cm(0.35)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Cm(0.5), bar_w, height - Cm(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    inner_left = left + bar_w + Cm(0.6)
    inner_width = Emu(int(width) - int(bar_w) - int(Cm(0.6)))

    qt = slide.shapes.add_textbox(inner_left, top, inner_width, Cm(2.5))
    tf_qt = qt.text_frame
    p_qt = tf_qt.paragraphs[0]
    p_qt.text = "\u201c"
    p_qt.font.size = Pt(72)
    p_qt.font.bold = True
    p_qt.font.color.rgb = accent_color
    p_qt.alignment = PP_ALIGN.LEFT

    quote_top = top + Cm(1.8)
    quote_h = Emu(int(height) - int(Cm(1.8)) - int(Cm(1.8)))
    qt_body = slide.shapes.add_textbox(inner_left, quote_top, inner_width, quote_h)
    tf_body = qt_body.text_frame
    tf_body.word_wrap = True
    tf_body.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_body = tf_body.paragraphs[0]
    p_body.text = quote
    p_body.font.size = Pt(24)
    p_body.font.bold = True
    p_body.font.color.rgb = text_color
    p_body.alignment = PP_ALIGN.LEFT

    attr_top = top + height - Cm(1.4)
    attr_box = slide.shapes.add_textbox(inner_left, attr_top, inner_width, Cm(1.2))
    tf_attr = attr_box.text_frame
    tf_attr.word_wrap = True
    p_attr = tf_attr.paragraphs[0]
    p_attr.text = attribution
    p_attr.font.size = Pt(13)
    p_attr.font.italic = True
    p_attr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p_attr.alignment = PP_ALIGN.LEFT


def add_progress_ring(
    slide,
    pct,
    label,
    value_text,
    cx,
    cy,
    diameter=Cm(3.5),
    ring_color=COLOR_ACCENT,
    bg_color=RGBColor(0xE0, 0xE0, 0xE0),
):
    """Circular KPI gauge: grey background donut + colored arc overlay.

    pct: float 0-100.
    label: str below the ring.
    value_text: str inside the ring, e.g. "74%".
    cx/cy: center position in EMU.
    """
    r = diameter
    left = cx - Emu(int(r) // 2)
    top = cy - Emu(int(r) // 2)
    hole = 0.55

    track = slide.shapes.add_shape(MSO_SHAPE.DONUT, left, top, r, r)
    track.adjustments[0] = hole
    track.fill.solid()
    track.fill.fore_color.rgb = bg_color
    track.line.fill.background()

    arc = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, left, top, r, r)
    end_angle = (pct / 100.0) * 360.0
    arc.fill.solid()
    arc.fill.fore_color.rgb = ring_color
    arc.line.fill.background()
    sp = arc._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is not None:
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is not None:
                for gd in list(avLst):
                    avLst.remove(gd)
                for name, val in [
                    ("adj1", 0),
                    ("adj2", int(end_angle * 60000)),
                    ("adj3", int(hole * 50000)),
                ]:
                    gd = etree.SubElement(avLst, qn("a:gd"))
                    gd.set("name", name)
                    gd.set("fmla", f"val {val}")

    val_box = slide.shapes.add_textbox(left, top, r, r)
    tf = val_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    lbl = slide.shapes.add_textbox(
        left - Cm(0.5), top + r + Cm(0.2), Emu(int(r) + int(Cm(1.0))), Cm(0.8)
    )
    tf2 = lbl.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER


def add_hexagon_grid(slide, items, left, top, width, hex_size=Cm(2.8), colors=None):
    """Honeycomb grid of hexagon tiles.

    items: list of str — labels for each hexagon.
    """
    h_w = hex_size
    h_h = Emu(int(int(h_w) * 0.87))
    gap = Cm(0.2)
    cols = int(int(width) / (int(h_w) + int(gap)))
    row_step = Emu(int(int(h_h) * 0.75) + int(gap))
    col_step = Emu(int(h_w) + int(gap))

    for idx, label in enumerate(items):
        row = idx // cols
        col = idx % cols
        offset = Emu(int(col_step) // 2) if row % 2 == 1 else Emu(0)
        hx_left = left + Emu(col * int(col_step)) + offset
        hx_top = top + Emu(row * int(row_step))
        color = colors[idx] if colors and idx < len(colors) else PALETTE[idx % len(PALETTE)]

        hx = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, hx_left, hx_top, h_w, h_h)
        hx.rotation = 90.0
        hx.fill.solid()
        hx.fill.fore_color.rgb = color
        hx.line.fill.background()
        tf = hx.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_venn(slide, circles, left, top, width, height):
    """2-3 overlapping semi-transparent circles with labels.

    circles: list of (label, color) tuples.
    """
    n = len(circles)
    d = min(Emu(int(width) // n + int(Cm(1.5))), height)
    overlap = int(int(d) * 0.30)
    total_w = n * int(d) - (n - 1) * overlap
    start_x = left + Emu((int(width) - total_w) // 2)
    cy = top + Emu((int(height) - int(d)) // 2)

    for i, (label, color) in enumerate(circles):
        cx = start_x + Emu(i * (int(d) - overlap))
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = color
        circle.line.width = Pt(2)
        srgbClr = circle.fill._xPr.find(
            f"{{{'http://schemas.openxmlformats.org/drawingml/2006/main'}}}solidFill/"
            f"{{{'http://schemas.openxmlformats.org/drawingml/2006/main'}}}srgbClr"
        )
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn("a:alpha"))
            alpha.set("val", "40000")

        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_funnel(slide, stages, left, top, width, height, colors=None):
    """Vertical funnel from stacked narrowing trapezoids.

    stages: list of (label, value_text) tuples, top-to-bottom.
    """
    n = len(stages)
    gap = Cm(0.15)
    stage_h = Emu(int((int(height) - int(gap) * (n - 1)) / n))
    if int(stage_h) < int(Cm(1.0)):
        raise ValueError(f"add_funnel: too many stages ({n}) for available height")
    center_x = left + Emu(int(width) // 2)

    for i, (label, value) in enumerate(stages):
        frac = 1.0 - (i * 0.7 / max(n - 1, 1))
        s_w = Emu(int(int(width) * frac))
        s_left = center_x - Emu(int(s_w) // 2)
        s_top = top + Emu(i * (int(stage_h) + int(gap)))
        color = colors[i] if colors and i < len(colors) else PALETTE[i % len(PALETTE)]

        trap = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, s_left, s_top, s_w, stage_h)
        trap.fill.solid()
        trap.fill.fore_color.rgb = color
        trap.line.fill.background()

        tf = trap.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = label
        run1.font.size = Pt(12)
        run1.font.bold = True
        run1.font.color.rgb = WHITE
        run2 = p.add_run()
        run2.text = f"  {value}"
        run2.font.size = Pt(12)
        run2.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_callout(
    slide,
    text,
    box_left,
    box_top,
    box_w,
    box_h,
    target_x,
    target_y,
    color=COLOR_PRIMARY,
):
    """Leader-line callout pointing from a label box to a target coordinate.

    Use to annotate screenshots, images, or charts.
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_w, box_h)
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_PRIMARY

    box_cx = box_left + Emu(int(box_w) // 2)
    start_y = box_top + box_h
    # Vertical leader line (thin rectangle instead of connector)
    line_h = Emu(int(target_y) - int(start_y))
    if int(line_h) > 0:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_cx - Pt(0.5), start_y, Pt(1), line_h)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

    dot_r = Cm(0.12)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, target_x - dot_r, target_y - dot_r, 2 * dot_r, 2 * dot_r
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
