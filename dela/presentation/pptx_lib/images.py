"""
pptx_lib.images — Image placement, overlays, and shadow boxes.
"""

from lxml import etree
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .helpers import send_to_back


def fill_placeholder_image(slide, idx, img_path):
    """Fill a picture placeholder with an image, with robust fallback.

    python-pptx's insert_picture() only works on genuine PicturePlaceholder
    objects.  After template manipulation (delete/re-add slides) the
    placeholder may lose its type, causing insert_picture() to silently fail
    or raise.

    Strategy:
    1. Try insert_picture() on the placeholder.
    2. If that fails, read the placeholder's position/size, remove it,
       and add the image as a regular shape via add_picture(), then
       send it to the back.

    Returns the picture shape.
    """
    ph = slide.placeholders[idx]
    left, top, width, height = ph.left, ph.top, ph.width, ph.height

    # Attempt native insert_picture (works on genuine PicturePlaceholders)
    if hasattr(ph, "insert_picture"):
        try:
            pic = ph.insert_picture(img_path)
            # Verify the image was actually embedded by checking for a blipFill
            blip = pic._element.find(".//" + qn("a:blip"))
            if blip is not None and blip.get(qn("r:embed")):
                return pic
        except Exception:
            pass

    # Fallback: remove placeholder, place image as a regular shape
    ph._element.getparent().remove(ph._element)

    # Preserve aspect ratio within the placeholder area
    img = Image.open(img_path)
    img_w, img_h = img.size
    img_ratio = img_w / img_h
    ph_ratio = width / height
    if img_ratio > ph_ratio:
        height = int(width / img_ratio)
        top = top + (height - height) // 2
    else:
        height = height
        width = int(height * img_ratio)
        left = left + (width - width) // 2
        top = top
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
    send_to_back(slide, pic)
    return pic


def add_image_preserved_ratio(
    slide, img_path, target_left, target_top, target_width, target_height
):
    """Add an image preserving its aspect ratio, centered within the target box.

    Returns the added picture shape.
    """
    img = Image.open(img_path)
    img_w, img_h = img.size
    img_ratio = img_w / img_h
    target_ratio = target_width / target_height
    if img_ratio > target_ratio:
        width = target_width
        height = int(width / img_ratio)
        left = target_left
        top = target_top + (target_height - height) // 2
    else:
        height = target_height
        width = int(height * img_ratio)
        left = target_left + (target_width - width) // 2
        top = target_top
    return slide.shapes.add_picture(img_path, left, top, width, height)


def add_text_shadow_box(
    slide, left, top, width, height, rgb=RGBColor(0x00, 0x00, 0x00), opacity_pct=55
):
    """Add a semi-transparent colored rectangle behind text on image slides.

    Placed between background image (z-pos 2) and text placeholders (z-pos 4+).
    opacity_pct: 0=fully transparent, 100=fully opaque.
    """
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = rgb
    alpha_val = int(opacity_pct * 1000)
    spPr = box._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        srgbClr = solidFill.find(qn("a:srgbClr"))
        if srgbClr is not None:
            alpha_el = etree.SubElement(srgbClr, qn("a:alpha"))
            alpha_el.set("val", str(alpha_val))
    send_to_back(slide, box)
    sp = box._element
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    spTree.insert(3, sp)
    return box


def add_background_image_with_overlay(slide, img_path, left, top, width, height, overlay_alpha=0.4):
    """Add a background image with a semi-transparent dark overlay.

    Good for text readability on image backgrounds.
    """
    add_image_preserved_ratio(slide, img_path, left, top, width, height)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    solidFill = overlay.fill._fill
    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is not None:
        alpha = etree.SubElement(srgbClr, qn("a:alpha"))
        alpha.set("val", str(int(overlay_alpha * 100000)))
    overlay.line.fill.background()
    return overlay
