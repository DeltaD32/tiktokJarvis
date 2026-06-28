"""
pptx_lib.credits — Attribution / image credits slide.
"""

from pptx.dml.color import RGBColor
from pptx.util import Pt


def add_credits_slide(prs, credits):
    """Add a final "Image Credits" slide listing all attributions.

    credits: list of (slide_num, description, author, source, license) tuples.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[7])  # Content | 1
    slide.placeholders[0].text = "Image Credits"
    tf = slide.placeholders[17].text_frame
    tf.clear()
    for i, (slide_num, desc, author, source, lic) in enumerate(credits):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"Slide {slide_num}: {desc} by {author} ({source}, {lic})"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x5A, 0x6B, 0x7A)
        p.space_before = Pt(4)
    return slide
