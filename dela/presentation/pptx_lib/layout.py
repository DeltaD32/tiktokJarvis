"""
pptx_lib.layout — Structured layout patterns: pillars, MECE tiles, pyramids,
sidebar, tabs, zigzag, numbered badges, key-value pairs.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt

from .cards import add_icon_circle
from .constants import (
    CARD_BG,
    CARD_BORDER,
    COLOR_ACCENT,
    COLOR_PRIMARY,
    DARK,
    PALETTE,
    WHITE,
    body_font_size,
)
from .grid import ContentGrid
from .helpers import send_to_back
from .text import add_accent_underline


def add_pillars(slide, pillars, left, top, total_width, total_height, colors=None):
    """Tall rounded rectangles side by side with colored header + light body.

    pillars: list of (title, [bullet1, bullet2, ...]) tuples.
    colors: optional list of RGBColor per pillar.
    """
    n = len(pillars)
    header_h = Cm(1.5)
    grid = ContentGrid(
        left=left,
        top=top,
        width=total_width,
        height=total_height,
        cols=n,
        rows=1,
        col_gap=Cm(0.4),
        row_gap=Cm(0),
    )

    for i, (title, bullets) in enumerate(pillars):
        cl, ct, cw, ch = grid.cell(i, 0)
        color = colors[i] if colors else COLOR_PRIMARY

        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, cw, header_h)
        header.adjustments[0] = 0.02
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        tf = header.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        body_top = ct + header_h + Emu(5000)
        body_h = ch - header_h - Emu(5000)
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, body_top, cw, body_h)
        body.adjustments[0] = 0.02
        body.fill.solid()
        body.fill.fore_color.rgb = CARD_BG
        body.line.color.rgb = CARD_BORDER
        body.line.width = Pt(1)
        send_to_back(slide, body)

        n_bullets = max(len(bullets), 1)
        body_h_cm = int(body_h) / int(Cm(1))
        fsize = body_font_size(body_h_cm / n_bullets)

        txt = slide.shapes.add_textbox(
            cl + Cm(0.4),
            body_top + Cm(0.4),
            Emu(int(cw) - int(Cm(0.8))),
            Emu(int(body_h) - int(Cm(0.8))),
        )
        tf2 = txt.text_frame
        tf2.word_wrap = True
        for j, item in enumerate(bullets):
            pp = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            pp.text = f"\u2022 {item}"
            pp.font.size = fsize
            pp.font.color.rgb = DARK
            pp.space_before = Pt(4)


def add_mece_tiles(slide, items, left, top, width, height, colors=None, show_numbers=True):
    """Equal-width vertical tiles for MECE breakdowns.

    items: list of (title, body_text) tuples — 2 to 5 tiles.
    colors: optional list of RGBColor; cycles PALETTE if omitted.
    show_numbers: add numbered icon circle at top of each tile.
    """
    n = len(items)
    header_h = Cm(1.6)
    grid = ContentGrid(
        left=left,
        top=top,
        width=width,
        height=height,
        cols=n,
        rows=1,
        col_gap=Cm(0.35),
        row_gap=Cm(0),
    )

    for i, (title, body) in enumerate(items):
        color = colors[i] if colors and i < len(colors) else PALETTE[i % len(PALETTE)]
        cl, ct, cw, ch = grid.cell(i, 0)

        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, cw, header_h)
        hdr.adjustments[0] = 0.02
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tf = hdr.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ph = tf.paragraphs[0]
        ph.text = title
        ph.font.size = Pt(16)
        ph.font.bold = True
        ph.font.color.rgb = WHITE
        ph.alignment = PP_ALIGN.CENTER

        body_top = ct + header_h + Emu(3600)
        body_h = ch - header_h - Emu(3600)
        body_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, body_top, cw, body_h)
        body_card.adjustments[0] = 0.02
        body_card.fill.solid()
        body_card.fill.fore_color.rgb = CARD_BG
        body_card.line.color.rgb = CARD_BORDER
        body_card.line.width = Pt(1)
        send_to_back(slide, body_card)

        body_h_cm = int(body_h) / int(Cm(1))
        fsize = body_font_size(body_h_cm)

        pad = Cm(0.4)
        txt = slide.shapes.add_textbox(
            cl + pad,
            body_top + pad,
            Emu(int(cw) - int(2 * pad)),
            Emu(int(body_h) - int(2 * pad)),
        )
        tf2 = txt.text_frame
        tf2.word_wrap = True
        pb = tf2.paragraphs[0]
        pb.text = body
        pb.font.size = fsize
        pb.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        if show_numbers:
            add_icon_circle(
                slide,
                cl + Emu(int(cw) // 2) - Emu(int(Cm(0.375))),
                ct - Cm(0.375),
                str(i + 1),
                color=color,
                diameter=Cm(0.75),
            )


def add_pyramid_slide_content(
    slide, headline, supports, left, top, width, height, accent_color=COLOR_ACCENT
):
    """Answer-first layout: bold headline bar + supporting argument cards.

    headline: str — the "So What?" conclusion.
    supports: list of (title, body) tuples — 2 to 4 arguments.
    """
    bar_h = Cm(1.8)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, bar_h)
    bar.adjustments[0] = 0.02
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    n = len(supports)
    card_top = top + bar_h + Cm(0.5)
    card_h = height - bar_h - Cm(0.5)
    grid = ContentGrid(
        left=left,
        top=card_top,
        width=width,
        height=card_h,
        cols=n,
        rows=1,
        col_gap=Cm(0.35),
        row_gap=Cm(0),
    )

    for i, (title, body) in enumerate(supports):
        cl, ct, cw, ch = grid.cell(i, 0)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, cw, ch)
        card.adjustments[0] = 0.02
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)
        send_to_back(slide, card)

        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cl, ct, Cm(0.15), ch)
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent_color
        strip.line.fill.background()
        send_to_back(slide, strip)

        body_area_h_cm = (int(ch) - int(Cm(1.6))) / int(Cm(1))
        fsize = body_font_size(body_area_h_cm)

        inner_left = cl + Cm(0.45)
        inner_w = Emu(int(cw) - int(Cm(0.6)))
        title_box = slide.shapes.add_textbox(inner_left, ct + Cm(0.3), inner_w, Cm(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        pt = tf2.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_PRIMARY

        body_box = slide.shapes.add_textbox(
            inner_left, ct + Cm(1.3), inner_w, Emu(int(ch) - int(Cm(1.6)))
        )
        tf3 = body_box.text_frame
        tf3.word_wrap = True
        pb = tf3.paragraphs[0]
        pb.text = body
        pb.font.size = fsize
        pb.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_sidebar_layout(
    slide,
    sidebar_items,
    main_title,
    main_text,
    left,
    top,
    width,
    height,
    sidebar_color=COLOR_PRIMARY,
):
    """Left sidebar with topic labels + wide main content area right.

    sidebar_items: list of str — short labels stacked in sidebar.
    """
    sidebar_w = Emu(int(int(width) * 0.20))
    main_left = left + sidebar_w + Cm(0.5)
    main_w = Emu(int(width) - int(sidebar_w) - int(Cm(0.5)))
    item_h = Emu(int(height) // max(len(sidebar_items), 1))
    if int(main_w) < int(Cm(5.0)):
        raise ValueError("add_sidebar_layout: main content area too narrow for given width")
    if int(item_h) < int(Cm(1.0)):
        raise ValueError(
            f"add_sidebar_layout: too many sidebar items ({len(sidebar_items)}) "
            "for available height"
        )

    sb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, sidebar_w, height)
    sb.adjustments[0] = 0.02
    sb.fill.solid()
    sb.fill.fore_color.rgb = sidebar_color
    sb.line.fill.background()

    for i, label in enumerate(sidebar_items):
        lbl = slide.shapes.add_textbox(
            left + Cm(0.3),
            top + Emu(i * int(item_h)),
            Emu(int(sidebar_w) - int(Cm(0.6))),
            item_h,
        )
        tf = lbl.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        if i < len(sidebar_items) - 1:
            sep_y = top + Emu((i + 1) * int(item_h))
            sep_w = Emu(int(sidebar_w) - int(Cm(1.0)))
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(0.5), sep_y, sep_w, Pt(0.5))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0x4A, 0x8A, 0xA0)
            sep.line.fill.background()

    mt = slide.shapes.add_textbox(main_left, top + Cm(0.2), main_w, Cm(1.0))
    tf_t = mt.text_frame
    tf_t.word_wrap = True
    pt = tf_t.paragraphs[0]
    pt.text = main_title
    pt.font.size = Pt(16)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_PRIMARY

    add_accent_underline(slide, main_left, top + Cm(1.4), Cm(5.0))

    main_body_h_cm = (int(height) - int(Cm(2.0))) / int(Cm(1))
    mb = slide.shapes.add_textbox(main_left, top + Cm(1.8), main_w, Emu(int(height) - int(Cm(2.0))))
    tf_m = mb.text_frame
    tf_m.word_wrap = True
    pm = tf_m.paragraphs[0]
    pm.text = main_text
    pm.font.size = body_font_size(main_body_h_cm)
    pm.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_tabbed_headers(
    slide,
    tabs,
    active_idx,
    left,
    top,
    width,
    active_color=COLOR_PRIMARY,
    inactive_color=RGBColor(0xDD, 0xDD, 0xDD),
):
    """Horizontal tab shapes; active tab highlighted, others muted.

    Returns (content_left, content_top, content_width) for content below.
    """
    n = len(tabs)
    tab_h = Cm(1.2)
    grid = ContentGrid(
        left=left,
        top=top,
        width=width,
        height=tab_h,
        cols=n,
        rows=1,
        col_gap=Cm(0.15),
        row_gap=Cm(0),
    )

    for i, label in enumerate(tabs):
        cl, ct, cw, ch = grid.cell(i, 0)
        is_active = i == active_idx
        color = active_color if is_active else inactive_color
        text_color = WHITE if is_active else RGBColor(0x66, 0x66, 0x66)

        tab = slide.shapes.add_shape(MSO_SHAPE.SNIP_2_SAME_RECTANGLE, cl, ct, cw, ch)
        tab.fill.solid()
        tab.fill.fore_color.rgb = color
        tab.line.fill.background()
        tf = tab.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = is_active
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    al, at, aw, ah = grid.cell(active_idx, 0)
    add_accent_underline(slide, al, at + ah, aw, color=active_color, thickness=Pt(4))

    return (left, top + tab_h + Cm(0.3), width)


def add_zigzag_blocks(
    slide,
    items,
    left,
    top,
    width,
    height,
    block_color=CARD_BG,
    accent_color=COLOR_ACCENT,
):
    """Alternating left-right content blocks with connecting lines.

    items: list of (title, text) tuples — 3 to 5 items.
    """
    n = len(items)
    gap_v = Cm(0.3)
    block_h = Emu(int((int(height) - int(gap_v) * (n - 1)) / n))
    block_w = Emu(int(int(width) * 0.55))
    if int(block_h) < int(Cm(1.0)):
        raise ValueError(f"add_zigzag_blocks: too many items ({n}) for available height")

    for i, (title, text) in enumerate(items):
        is_left = i % 2 == 0
        b_left = left if is_left else (left + Emu(int(width) - int(block_w)))
        b_top = top + Emu(i * (int(block_h) + int(gap_v)))

        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, b_left, b_top, block_w, block_h)
        block.adjustments[0] = 0.02
        block.fill.solid()
        block.fill.fore_color.rgb = block_color
        block.line.color.rgb = CARD_BORDER
        block.line.width = Pt(1)
        send_to_back(slide, block)

        strip_left = (b_left + block_w - Cm(0.12)) if is_left else b_left
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, strip_left, b_top, Cm(0.12), block_h)
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent_color
        strip.line.fill.background()

        pad = Cm(0.4)
        inner_left = b_left + pad
        inner_w = Emu(int(block_w) - int(2 * pad))
        tb = slide.shapes.add_textbox(inner_left, b_top + Cm(0.15), inner_w, Cm(0.6))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        pt = tf_t.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(15)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_PRIMARY

        body_area_h_cm = (int(block_h) - int(Cm(1.1))) / int(Cm(1))
        bb = slide.shapes.add_textbox(
            inner_left, b_top + Cm(0.9), inner_w, Emu(int(block_h) - int(Cm(1.1)))
        )
        tf_b = bb.text_frame
        tf_b.word_wrap = True
        pb = tf_b.paragraphs[0]
        pb.text = text
        pb.font.size = body_font_size(body_area_h_cm)
        pb.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        if i < n - 1:
            curr_edge_x = (b_left + block_w) if is_left else b_left
            next_is_left = not is_left
            next_b_left = left if next_is_left else (left + Emu(int(width) - int(block_w)))
            next_edge_x = (next_b_left + block_w) if next_is_left else next_b_left
            # Use thin rectangle instead of connector to avoid repair issues
            mid_y = b_top + Emu(int(block_h) // 2)
            min_x = min(int(curr_edge_x), int(next_edge_x))
            conn_w = abs(int(curr_edge_x) - int(next_edge_x))
            if conn_w > 0:
                conn = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(min_x), mid_y, Emu(conn_w), Pt(1)
                )
                conn.fill.solid()
                conn.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                conn.line.fill.background()
                send_to_back(slide, conn)


def add_numbered_badge_list(slide, items, left, top, width, height, badge_color=COLOR_PRIMARY):
    """Vertical list with numbered badge circles and text.

    items: list of (title, description) tuples — numbered automatically.
    """
    n = len(items)
    badge_d = Cm(1.2)
    text_left = left + badge_d + Cm(0.5)
    text_w = Emu(int(width) - int(badge_d) - int(Cm(0.5)))
    grid = ContentGrid(
        left=left,
        top=top,
        width=width,
        height=height,
        cols=1,
        rows=n,
        col_gap=Cm(0),
        row_gap=Cm(0),
    )

    for i, (title, desc) in enumerate(items):
        cl, ct, cw, ch = grid.cell(0, i)
        badge_top = ct + Emu((int(ch) - int(badge_d)) // 2)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, badge_top, badge_d, badge_d)
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        pb = tf_b.paragraphs[0]
        pb.text = str(i + 1)
        pb.font.size = Pt(18)
        pb.font.bold = True
        pb.font.color.rgb = WHITE
        pb.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(text_left, ct + Cm(0.2), text_w, Cm(0.7))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        pt = tf_t.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_PRIMARY

        desc_area_h_cm = (int(ch) - int(Cm(1.2))) / int(Cm(1))
        desc_box = slide.shapes.add_textbox(
            text_left, ct + Cm(1.0), text_w, Emu(int(ch) - int(Cm(1.2)))
        )
        tf_d = desc_box.text_frame
        tf_d.word_wrap = True
        pd = tf_d.paragraphs[0]
        pd.text = desc
        pd.font.size = body_font_size(desc_area_h_cm)
        pd.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        if i < n - 1:
            # Use a thin rectangle instead of add_connector to avoid repair issues
            line_x = left + Emu(int(badge_d) // 2) - Emu(int(Pt(0.5)))
            _, next_ct, _, next_ch = grid.cell(0, i + 1)
            next_badge_top = next_ct + Emu((int(next_ch) - int(badge_d)) // 2)
            line_top = badge_top + badge_d
            line_h = Emu(int(next_badge_top) - int(line_top))
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x, line_top, Pt(1), line_h)
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            line.line.fill.background()


def add_key_value_pairs(slide, pairs, left, top, width, height, accent_color=COLOR_ACCENT):
    """Horizontal key-value rows with divider lines and left accent bar.

    pairs: list of (key, value) str tuples.
    """
    n = len(pairs)
    key_w = Emu(int(int(width) * 0.35))
    val_left = left + key_w + Cm(0.5)
    val_w = Emu(int(width) - int(key_w) - int(Cm(0.5)))
    grid = ContentGrid(
        left=left,
        top=top,
        width=width,
        height=height,
        cols=1,
        rows=n,
        col_gap=Cm(0),
        row_gap=Cm(0),
    )

    for i, (key, value) in enumerate(pairs):
        cl, ct, cw, ch = grid.cell(0, i)

        kb = slide.shapes.add_textbox(left, ct, key_w, ch)
        tf_k = kb.text_frame
        tf_k.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_k.word_wrap = True
        pk = tf_k.paragraphs[0]
        row_h_cm = int(ch) / int(Cm(1))
        pk.text = key
        pk.font.size = body_font_size(row_h_cm)
        pk.font.bold = True
        pk.font.color.rgb = COLOR_PRIMARY
        pk.alignment = PP_ALIGN.RIGHT

        vb = slide.shapes.add_textbox(val_left, ct, val_w, ch)
        tf_v = vb.text_frame
        tf_v.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_v.word_wrap = True
        pv = tf_v.paragraphs[0]
        row_h_cm = int(ch) / int(Cm(1))
        pv.text = value
        pv.font.size = body_font_size(row_h_cm)
        pv.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        if i < n - 1:
            line_y = ct + ch
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, line_y, width, Pt(0.75))
            ln.fill.solid()
            ln.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            ln.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Cm(0.2), top, Cm(0.12), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
