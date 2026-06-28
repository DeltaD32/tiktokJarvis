"""
pptx_lib.helpers — Low-level slide manipulation helpers.

Z-order, placeholder cleanup, bullet suppression, speaker notes,
department badge geometry preservation.
"""

from copy import deepcopy

from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Cm

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def send_to_back(slide, shape):
    """Move shape to bottom of z-order (drawn first = behind everything)."""
    sp = shape._element
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    spTree.insert(2, sp)


def bring_to_front(slide, shape):
    """Move shape to top of z-order (drawn last = in front of everything)."""
    sp = shape._element
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    spTree.append(sp)


def remove_empty_placeholders(slide):
    """Remove unfilled content placeholders (idx != 0) from a slide.

    Prevents 'Click to edit text' ghost placeholders over custom shapes.
    Preserves:
    - Title placeholder (idx 0)
    - PicturePlaceholder objects (has insert_picture attribute)
    - Placeholders containing an embedded image (a:blip or p:pic in XML),
      which happens after insert_picture() transforms the element
    """
    to_remove = []
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        # Keep if it's a picture placeholder (before insertion)
        if hasattr(ph, "insert_picture"):
            continue
        # Keep if it contains an embedded image (after insertion —
        # insert_picture transforms the XML, losing the Python attribute)
        el = ph._element
        if (
            el.find(".//" + qn("a:blip")) is not None
            or el.tag == qn("p:pic")
            or el.find(qn("p:pic")) is not None
        ):
            continue
        if ph.has_text_frame:
            if not ph.text_frame.text.strip():
                to_remove.append(ph)
        else:
            to_remove.append(ph)
    for ph in to_remove:
        ph._element.getparent().remove(ph._element)


def suppress_bullet(p):
    """Remove inherited template bullet from a paragraph.

    Use on all header/plain-text paragraphs in placeholders to prevent
    scissors/symbol glyph rendering from the BMW template.
    """
    pPr = p._p.get_or_add_pPr()
    for tag in ["a:buNone", "a:buChar", "a:buFont", "a:buClr", "a:buSzPct"]:
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    etree.SubElement(pPr, qn("a:buNone"))


def set_ph_margins(tf, left=Cm(0.55), top=Cm(0.4), right=Cm(0.4), bottom=Cm(0.3)):
    """Set text frame margins on a placeholder.

    Prevents text touching card edges. Defaults tuned for cards with accent strips.
    """
    tf.margin_left = left
    tf.margin_top = top
    tf.margin_right = right
    tf.margin_bottom = bottom


def set_notes(slide, text):
    """Set speaker notes for a slide."""
    slide.notes_slide.notes_text_frame.text = text


def fix_dept_placeholder(slide):
    """Copy layout-level spPr onto the slide-level department placeholder.

    python-pptx creates slide placeholders with an empty ``<p:spPr/>``.
    PowerPoint inherits geometry from the layout in this case, but
    LibreOffice (used for self-review rendering) does not — it falls back
    to a plain rectangle. The BMW template defines a custom trapezoid
    shape for placeholder 22. This function copies the full ``spPr``
    (geometry, fill, line, effects) from the layout so rendering is
    correct in both PowerPoint and LibreOffice.

    Call this **after** setting ``slide.placeholders[22].text``.
    Safe to call on any slide — silently returns if ph 22 is missing.
    """
    try:
        ph = slide.placeholders[22]
    except (KeyError, IndexError):
        return

    layout = slide.slide_layout
    layout_spPr = None
    for layout_ph in layout.placeholders:
        if layout_ph.placeholder_format.idx == 22:
            layout_spPr = layout_ph._element.find(f"{{{_P_NS}}}spPr")
            break
    if layout_spPr is None:
        return

    slide_sp = ph._element
    old_spPr = slide_sp.find(f"{{{_P_NS}}}spPr")
    new_spPr = deepcopy(layout_spPr)
    if old_spPr is not None:
        slide_sp.replace(old_spPr, new_spPr)
    else:
        # Insert after nvSpPr (first child)
        slide_sp.insert(1, new_spPr)
