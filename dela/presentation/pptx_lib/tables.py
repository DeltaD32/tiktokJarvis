"""
pptx_lib.tables — Styled tables and comparison tables.
"""

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .constants import COLOR_PRIMARY


def add_styled_table(
    slide, headers, rows, left, top, width, height, header_color=COLOR_PRIMARY, col_widths=None
):
    """Add a table with a styled header row.

    headers: list of str.
    rows: list of list of str (one inner list per row).
    col_widths: optional list of float fractions summing to 1.0.
    """
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    if col_widths:
        for c, frac in enumerate(col_widths):
            tbl.columns[c].width = Emu(int(int(width) * frac))

    # Header row
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = hdr
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    ROW_BG_ODD = RGBColor(0xF5, 0xF8, 0xFA)
    ROW_BG_EVEN = RGBColor(0xFF, 0xFF, 0xFF)
    for r, row_data in enumerate(rows):
        row_bg = ROW_BG_ODD if r % 2 == 0 else ROW_BG_EVEN
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            tf = cell.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.alignment = PP_ALIGN.CENTER

    return tbl_shape


def add_comparison_table(slide, criteria, options, left, top, width, height, option_colors=None):
    """Side-by-side option comparison with styled header and alternating rows.

    criteria: list of str — row labels (left column).
    options:  list of (option_name, [cell_text, ...]) tuples.
    option_colors: optional list of RGBColor per option column header.
    """
    n_cols = 1 + len(options)
    n_rows = 1 + len(criteria)
    if option_colors is None:
        option_colors = [COLOR_PRIMARY] * len(options)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    crit_w = int(int(width) * 0.25)
    opt_w = int((int(width) - crit_w) / len(options))
    tbl.columns[0].width = Emu(crit_w)
    for c in range(1, n_cols):
        tbl.columns[c].width = Emu(opt_w)

    # Top-left corner
    corner = tbl.cell(0, 0)
    corner.fill.solid()
    corner.fill.fore_color.rgb = COLOR_PRIMARY

    # Option headers
    for c_idx, (opt_name, _) in enumerate(options):
        cell = tbl.cell(0, c_idx + 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = option_colors[c_idx]
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = opt_name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    ROW_BG_ODD = RGBColor(0xF5, 0xF8, 0xFA)
    ROW_BG_EVEN = RGBColor(0xFF, 0xFF, 0xFF)
    for r_idx, criterion in enumerate(criteria):
        row_num = r_idx + 1
        row_bg = ROW_BG_ODD if r_idx % 2 == 0 else ROW_BG_EVEN

        crit_cell = tbl.cell(row_num, 0)
        crit_cell.fill.solid()
        crit_cell.fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF3)
        tf_c = crit_cell.text_frame
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
        pc = tf_c.paragraphs[0]
        pc.text = criterion
        pc.font.size = Pt(11)
        pc.font.bold = True
        pc.font.color.rgb = COLOR_PRIMARY

        for c_idx, (_, cell_values) in enumerate(options):
            cell = tbl.cell(row_num, c_idx + 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            tf = cell.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = cell_values[r_idx]
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.alignment = PP_ALIGN.CENTER

    return tbl_shape
