"""
pptx_lib.shapes — Chevron chains, arrow-bar steps, divider lines.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt

from .cards import add_bottom_banner, add_icon_circle
from .constants import (
    COLOR_ACCENT,
    COLOR_PRIMARY,
    FREE_HEIGHT,
    FREE_LEFT,
    FREE_TOP,
    FREE_WIDTH,
    WHITE,
)


def add_chevron_chain(
    slide,
    items,
    left=None,
    top=None,
    width=None,
    available_height=None,
    chev_height=Cm(5.0),
    overlap=Cm(1.5),
    badge_diameter=Cm(1.2),
    chev_color=COLOR_PRIMARY,
    badge_color=COLOR_ACCENT,
    kicker_text=None,
    kicker_height=Cm(2.0),
):
    """Horizontal chain of chevron arrows with numbered badge circles.

    items: list of (number_str, title, subtitle) tuples.
    First shape is PENTAGON (flat left), rest are CHEVRON.
    Badge circles float above each chevron at 40% horizontal offset.
    Entire group is vertically centered in the available zone.

    If kicker_text is provided, a bottom banner is added below.

    Returns list of (chevron_shape, badge_shape) tuples.
    """
    left = left or (FREE_LEFT + Cm(0.2))
    top = top or FREE_TOP
    width = width or (FREE_WIDTH - Cm(0.4))
    available_height = available_height or FREE_HEIGHT

    BADGE_D = badge_diameter
    CHEV_H = chev_height
    KICKER_H = kicker_height if kicker_text else Cm(0)

    # Vertical centering
    group_h = BADGE_D / 2 + CHEV_H
    zone_h = available_height - KICKER_H - Cm(0.4)
    badge_top = top + (zone_h - group_h) / 2
    chev_top = badge_top + BADGE_D / 2

    n = len(items)
    chev_w = (width + (n - 1) * overlap) / n
    if int(chev_w) < int(Cm(2.5)):
        raise ValueError(f"add_chevron_chain: too many items ({n}) for available width")

    results = []
    for i, (num, label, sublabel) in enumerate(items):
        cleft = left + i * (chev_w - overlap)

        shape_type = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        chev = slide.shapes.add_shape(shape_type, cleft, chev_top, chev_w, CHEV_H)
        chev.fill.solid()
        chev.fill.fore_color.rgb = chev_color
        chev.line.fill.background()

        # Badge at 40% of chev width to clear arrow tip
        badge_cx = cleft + chev_w * 0.40
        badge_left = badge_cx - BADGE_D / 2
        badge = add_icon_circle(
            slide, badge_left, badge_top, num, color=badge_color, diameter=BADGE_D
        )

        # Text padding per shape geometry
        if i == 0:
            lpad, rpad = Cm(0.6), Cm(1.2)
        elif i == n - 1:
            lpad, rpad = Cm(1.3), Cm(0.4)
        else:
            lpad, rpad = Cm(1.3), Cm(1.2)

        tb = slide.shapes.add_textbox(
            cleft + lpad,
            chev_top + BADGE_D * 0.6,
            chev_w - lpad - rpad,
            CHEV_H - BADGE_D * 0.6,
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = sublabel
        p2.font.size = Pt(13)
        p2.font.bold = False
        p2.font.color.rgb = RGBColor(0xCC, 0xE6, 0xF0)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(5)

        results.append((chev, badge))

    if kicker_text:
        add_bottom_banner(
            slide,
            kicker_text,
            left,
            top + available_height - KICKER_H,
            width,
            bg_color=RGBColor(0x1A, 0x3A, 0x48),
            font_size=Pt(20),
        )

    return results


def add_arrow_bar_steps(slide, items, bar_left, bar_top, bar_width, arrow_height, colors=None):
    """Down-arrows hanging from a horizontal bar.

    items: list of (number, label, description) tuples.
    """
    bar_h = Cm(0.4)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, bar_width, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    bar.line.fill.background()

    n = len(items)
    spacing = int(int(bar_width) / n)
    arrow_w = Cm(2.5)

    for i, (num, label, desc) in enumerate(items):
        center_x = int(bar_left) + int(spacing * (i + 0.5))
        arrow_left = Emu(center_x - int(arrow_w) // 2)
        a_top = bar_top + bar_h
        color = colors[i] if colors else COLOR_PRIMARY

        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, arrow_left, a_top, arrow_w, arrow_height
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
        tf = arrow.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        lbl_top = a_top + arrow_height + Cm(0.3)
        lbl = slide.shapes.add_textbox(Emu(center_x - int(Cm(3)) // 2), lbl_top, Cm(3), Cm(0.8))
        tf2 = lbl.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

        if desc:
            desc_box = slide.shapes.add_textbox(
                Emu(center_x - int(Cm(3.5)) // 2), lbl_top + Cm(0.9), Cm(3.5), Cm(2)
            )
            tf3 = desc_box.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = desc
            p3.font.size = Pt(10)
            p3.font.color.rgb = COLOR_PRIMARY
            p3.alignment = PP_ALIGN.CENTER


def add_divider_line(slide, left, top, width, color=RGBColor(0xCC, 0xCC, 0xCC)):
    """Thin horizontal line to visually separate content sections.
    Uses a thin rectangle instead of add_connector to avoid PowerPoint repair dialogs.
    """
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_gap_chevron(
    slide,
    left_box_left,
    left_box_width,
    right_box_left,
    center_y,
    height=Cm(2.0),
    color=COLOR_ACCENT,
    side_margin=Cm(0.2),
    max_width=Cm(2.2),
):
    """Add a chevron centered in the gap between two boxes without overlap.

    The chevron width is clamped to the available gap minus side margins.
    Raises ValueError if the gap is too small for a readable chevron.
    """
    left_right = left_box_left + left_box_width
    gap = right_box_left - left_right
    usable = gap - 2 * side_margin
    min_width = Cm(0.8)
    if usable < min_width:
        raise ValueError(f"Gap too small for chevron: gap={int(gap)} emu, usable={int(usable)} emu")

    chev_w = min(max_width, usable)
    chev_left = left_right + (gap - chev_w) / 2
    chev_top = center_y - height / 2

    chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, chev_left, chev_top, chev_w, height)
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = color
    chevron.line.fill.background()
    return chevron
