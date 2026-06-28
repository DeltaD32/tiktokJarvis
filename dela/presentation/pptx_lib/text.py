"""
pptx_lib.text — Text boxes, bullet lists, and accent underlines.
"""

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from .constants import COLOR_ACCENT, DARK


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=Pt(14),
    bold=False,
    color=DARK,
    alignment=PP_ALIGN.LEFT,
    word_wrap=True,
    vertical_anchor=None,
    anchor=None,
):
    """Add a textbox with basic formatting. Returns (textbox_shape, text_frame).

    ``anchor`` is accepted as an alias for ``vertical_anchor`` (common LLM typo).
    """
    if anchor is not None and vertical_anchor is None:
        vertical_anchor = anchor
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if vertical_anchor:
        tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox, tf


def add_bullet_paragraph(tf, text, font_size=Pt(14), color=DARK, bold=False):
    """Add a bullet paragraph to a text frame.

    Uses XML to create a reliable Unicode bullet (•) with Arial font.
    Removes any inherited template bullet elements first.
    """
    p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(4)
    pPr = p._p.get_or_add_pPr()
    for tag in ["a:buNone", "a:buChar", "a:buFont"]:
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "\u2022")
    return p


def add_accent_underline(slide, left, top, width, color=COLOR_ACCENT, thickness=Pt(3)):
    """Thin colored bar below a heading.

    Place directly below a heading textbox for a lightweight visual anchor.
    """
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line
