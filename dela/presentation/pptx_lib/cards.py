"""
pptx_lib.cards — Card backgrounds, accent strips, icon circles, bottom banners.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

from .constants import (
    CARD_BG,
    CARD_BORDER,
    COLOR_ACCENT,
    COLOR_PRIMARY,
    WHITE,
)
from .helpers import send_to_back


def add_card_with_accent(slide, left, top, width, height, accent_color=COLOR_ACCENT):
    """Add a rounded rectangle card with a colored accent strip on its left edge.

    The accent strip sits directly LEFT of the card (left - 0.15cm), inset
    top/bottom by 0.25cm to avoid protruding past rounded corners.
    The strip stays in front of the card background (do not send_to_back).

    Returns (card_shape, strip_shape).
    """
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.02
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1)
    send_to_back(slide, card)

    inset = Cm(0.25)
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left - Cm(0.15), top + inset, Cm(0.15), height - 2 * inset
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent_color
    strip.line.fill.background()
    return card, strip


def add_icon_circle(slide, left, top, text, color=COLOR_ACCENT, diameter=Cm(0.85)):
    """Add a small colored circle with centered text.

    Use for step numbering, pillar initials, or category icons.
    Returns the icon shape.
    """
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()
    tf = icon.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return icon


def add_bottom_banner(
    slide,
    text,
    left,
    top,
    width,
    bg_color=COLOR_PRIMARY,
    text_color=WHITE,
    font_size=Pt(20),
    banner_height=Cm(2.1),
):
    """Full-width colored banner with centered bold text.

    Standard banner height: 2.1 cm.
    Position: top = FREE_TOP + card_h + Cm(0.3).
    Size cards as: card_h = FREE_HEIGHT - Cm(2.4) to leave room.

    Returns the banner shape.
    """
    BANNER_H = banner_height
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, BANNER_H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_color
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.08)
    tf.margin_bottom = Cm(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size if int(font_size) >= int(Pt(18)) else Pt(18)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return banner


def add_takeaway_box(
    slide,
    text,
    left,
    top,
    width,
    height=Cm(2.2),
    border_color=COLOR_PRIMARY,
    text_color=COLOR_PRIMARY,
    font_size=Pt(20),
):
    """Light boxed takeaway with strong typography.

    Use this when a full-width dark banner would be too heavy, but the bottom
    conclusion still needs to read as primary content.
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.02
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFA)
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.5)
    tf.margin_right = Cm(0.5)
    tf.margin_top = Cm(0.12)
    tf.margin_bottom = Cm(0.12)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size if int(font_size) >= int(Pt(18)) else Pt(18)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return box
